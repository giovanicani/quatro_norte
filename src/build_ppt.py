# Gera docs/Apresentacao_QuatroNorte.pptx — deck academico completo (perguntas 1-11)
# Y = custo de manutencao interno por km (CAD, deflacionado por CPI Canada, base dez/2025)
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PROJECT_ROOT = Path.cwd()
if PROJECT_ROOT.name in ("notebooks", "src"):
    PROJECT_ROOT = PROJECT_ROOT.parent
TABLES = PROJECT_ROOT / "reports" / "tables"
FIG = PROJECT_ROOT / "reports" / "figures"
FIG_EDA = FIG / "eda"
OUT = PROJECT_ROOT / "docs" / "Apresentacao_QuatroNorte.pptx"

# ---------- paleta ----------
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
DARK = RGBColor(0x2B, 0x2B, 0x2B)
ORANGE = RGBColor(0xC4, 0x7A, 0x3B)
GREY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT = RGBColor(0xF5, 0xF2, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x3B, 0x6E, 0xA5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return s


def tb(slide, x, y, w, h, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT,
       font="Calibri", line_spacing=1.0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run(); run.text = ln
        f = run.font; f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.name = font
    return box


def bullets(slide, x, y, w, h, items, size=16, color=DARK, gap=6):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, (txt, lvl, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(gap)
        run = p.add_run(); run.text = ("• " if lvl == 0 else "– ") + txt
        f = run.font; f.size = Pt(size - 2 * lvl); f.color.rgb = color; f.bold = bold; f.name = "Calibri"
    return box


def header(slide, kicker, title, color=NAVY):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()
    bar.shadow.inherit = False
    tb(slide, Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.35), kicker.upper(),
       size=12, color=ORANGE, bold=True)
    tb(slide, Inches(0.55), Inches(0.62), Inches(12.2), Inches(0.85), title,
       size=30, color=color, bold=True)


def pic_fit(slide, path, x, y, max_w, max_h):
    from PIL import Image
    with Image.open(path) as im:
        w, h = im.size
    ratio = min(max_w / w, max_h / h)
    pw, ph = int(w * ratio), int(h * ratio)
    slide.shapes.add_picture(str(path), x + int((max_w - pw) / 2), y + int((max_h - ph) / 2), pw, ph)


def divider(kicker, title, sub=""):
    s = add_slide(NAVY)
    tb(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(0.5), kicker.upper(),
       size=16, color=ORANGE, bold=True)
    tb(s, Inches(0.8), Inches(3.1), Inches(11.7), Inches(1.4), title, size=40, color=WHITE, bold=True)
    if sub:
        tb(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.2), sub, size=18, color=LIGHT)
    return s


fmt = lambda v, d=3: f"{v:,.{d}f}".replace(",", "X").replace(".", ",").replace("X", ".")

# ---------- dados dinamicos ----------
stats = pd.read_csv(TABLES / "03c_stats_ppt.csv", index_col=0)["valor"]
desc = pd.read_csv(TABLES / "03b_estatisticas_descritivas.csv")
corr = pd.read_csv(TABLES / "03b_correlacao_com_y.csv")
eta = pd.read_csv(TABLES / "03b_eta_categoricas.csv")
vif = pd.read_csv(TABLES / "03b_vif.csv")
freq = pd.read_csv(TABLES / "03b_frequencia_categorias.csv")
ystat = pd.read_csv(TABLES / "03b_y_por_categoria.csv")
evo_y = pd.read_csv(TABLES / "03c_evolucao_y_detalhe.csv")
defl_anual = pd.read_csv(TABLES / "04_comparacao_nominal_deflacionado.csv")

try:
    metricas = pd.read_csv(TABLES / "05_metricas_modelos.csv")
except FileNotFoundError:
    metricas = None
try:
    perm = pd.read_csv(TABLES / "05_importancia_permutacao_random_forest.csv")
except FileNotFoundError:
    perm = None
try:
    modelo_rec = pd.read_csv(TABLES / "05_modelo_recomendado.csv")
except FileNotFoundError:
    modelo_rec = None

TARGET = "custo_manutencao_interno_por_km_deflacionado"

# ============================================================ CAPA
s = add_slide(NAVY)
tb(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5), "MBA · PROJETO APLICADO · CIÊNCIA DE DADOS", size=15, color=ORANGE, bold=True)
tb(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.9),
   "Previsão de Custos de Manutenção\nde Carretas", size=44, color=WHITE, bold=True)
tb(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.9),
   "Modelagem preditiva do custo de manutenção interno por quilômetro\nda frota própria de carretas · janela 2020–2025", size=19, color=LIGHT)
bar = s.shapes.add_shape(1, Inches(0.9), Inches(5.5), Inches(2.6), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background(); bar.shadow.inherit = False
tb(s, Inches(0.9), Inches(5.8), Inches(11.5), Inches(0.8),
   "Quatro Norte Consulting  ·  FGV", size=16, color=LIGHT, bold=True)

# ============================================================ EMPRESA
s = add_slide()
header(s, "A empresa", "Quatro Norte Consulting")
bullets(s, Inches(0.55), Inches(1.7), Inches(6.6), Inches(5.0), [
    ("Consultoria de dados que atende uma operação de leasing/rental de carretas no Canadá", 0, True),
    ("Frota própria (cus_id_owner = 4): ~10,4 mil carretas ativas na janela analisada", 0, False),
    ("Carretas secas e refrigeradas (reefer) de até 53 pés, operando em províncias canadenses (ON, QC, ...)", 0, False),
    ("Contratos de leasing e rental com diferentes coberturas de manutenção (MAINT / NET / MIX)", 0, False),
    ("Manutenção própria: rede de oficinas registra ordens de serviço (OS) com mão de obra e peças", 0, False),
    ("Dor de negócio: parte relevante do custo de manutenção é absorvida pela empresa (custo interno, charge_flag = 'I') e precisa ser prevista para orçamento e precificação de contratos", 0, True),
])
card = s.shapes.add_shape(1, Inches(7.5), Inches(1.9), Inches(5.2), Inches(4.4))
card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.color.rgb = ORANGE; card.shadow.inherit = False
tb(s, Inches(7.8), Inches(2.1), Inches(4.6), Inches(0.4), "A OPERAÇÃO EM NÚMEROS", size=13, color=ORANGE, bold=True)
kpis = [
    (f"{int(stats['carretas']):,}".replace(",", "."), "carretas na frota própria"),
    ("238.818", "ordens de serviço (2020–2025)"),
    ("CAD 79,0 mi", "custo interno de manutenção no período"),
    ("6 anos", "de histórico: OS, odômetro, contratos, GPS"),
]
yy = 2.6
for big, small in kpis:
    tb(s, Inches(7.8), Inches(yy), Inches(4.6), Inches(0.5), big, size=24, color=NAVY, bold=True)
    tb(s, Inches(7.8), Inches(yy + 0.42), Inches(4.6), Inches(0.35), small, size=12, color=GREY)
    yy += 0.95

# ============================================================ AGENDA
s = add_slide()
header(s, "Roteiro", "Agenda da apresentação")
bullets(s, Inches(0.55), Inches(1.8), Inches(12.2), Inches(5.2), [
    ("Bloco 1 — Contexto, problema, objetivos, hipóteses e artigos científicos", 0, True),
    ("Contexto da pesquisa · pergunta do problema · objetivo geral e específicos · hipóteses · 4 artigos", 1, False),
    ("Bloco 2 — Base de dados, feature engineering, EDA e técnicas", 0, True),
    ("Modelo de dados (7 tabelas) · variável-alvo · deflação CPI Canadá · features derivadas · análise exploratória variável a variável · técnicas estatísticas e de machine learning", 1, False),
    ("Bloco 3 — Referencial teórico e metodologia", 0, True),
    ("Síntese do referencial · metodologia passo a passo (pipeline em 7 notebooks reprodutíveis)", 1, False),
])

# ============================================================ BLOCO 1
divider("Bloco 1", "Contexto, problema, objetivos e hipóteses",
        "Perguntas 1–5 da agenda da disciplina")

# contexto
s = add_slide()
header(s, "Pergunta 1", "Onde o projeto será realizado (contexto)")
bullets(s, Inches(0.55), Inches(1.8), Inches(12.2), Inches(5.0), [
    ("Empresa de leasing/rental de carretas no Canadá, atendida pela Quatro Norte Consulting", 0, True),
    ("Foco na frota própria (cus_id_owner = 4), com carretas ativas e leituras de quilometragem válidas entre 2020-01-01 e 2025-12-31", 0, False),
    ("A empresa executa manutenção preventiva e corretiva; parte do custo é repassada ao cliente e parte é absorvida pela operação (custo interno, charge_flag = 'I')", 0, False),
    ("Todos os custos são registrados em dólares canadenses (CAD)", 0, False),
    ("Dados extraídos do ERP de manutenção em modelo estrela: cadastro das carretas, ordens de serviço, mão de obra, peças, odômetro, contratos e telemetria GPS", 0, False),
    ("Decisões que o projeto apoia: orçamento anual de manutenção, precificação de contratos e priorização de frota", 0, True),
])

# pergunta problema
s = add_slide()
header(s, "Pergunta 2", "Pergunta do problema")
card = s.shapes.add_shape(1, Inches(1.2), Inches(2.3), Inches(10.9), Inches(2.2))
card.fill.solid(); card.fill.fore_color.rgb = NAVY; card.line.fill.background(); card.shadow.inherit = False
tb(s, Inches(1.7), Inches(2.75), Inches(9.9), Inches(1.5),
   "Quais são os fatores que mais influenciam o custo de manutenção\ninterno das carretas — e como prever esse custo por km futuro\ncom base nos dados históricos?",
   size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
bullets(s, Inches(1.2), Inches(5.0), Inches(10.9), Inches(1.8), [
    ("Variável de interesse: custo de manutenção interno por quilômetro (CAD/km), no grão carreta × mês", 0, True),
    ("Custo interno = absorvido pela empresa (charge_flag = 'I'), incluindo manutenções preventivas e corretivas", 0, False),
])

# objetivos
s = add_slide()
header(s, "Pergunta 3", "Objetivo geral e objetivos específicos")
tb(s, Inches(0.55), Inches(1.75), Inches(12.2), Inches(0.4), "OBJETIVO GERAL", size=14, color=ORANGE, bold=True)
tb(s, Inches(0.55), Inches(2.15), Inches(12.2), Inches(1.0),
   "Analisar os dados históricos de manutenção das carretas para identificar os principais fatores que influenciam o custo interno por km e desenvolver um modelo preditivo capaz de estimar custos futuros.",
   size=18, color=DARK)
tb(s, Inches(0.55), Inches(3.35), Inches(12.2), Inches(0.4), "OBJETIVOS ESPECÍFICOS", size=14, color=ORANGE, bold=True)
bullets(s, Inches(0.55), Inches(3.75), Inches(12.2), Inches(3.3), [
    ("Coletar, consolidar e organizar os dados históricos de manutenção das carretas", 0, False),
    ("Realizar análise exploratória dos dados para identificar padrões, tendências e variáveis relevantes", 0, False),
    ("Investigar a relação entre características dos contratos de leasing e os custos de manutenção", 0, False),
    ("Identificar os principais fatores associados aos custos de manutenção", 0, False),
    ("Desenvolver e avaliar modelos preditivos para estimar os custos futuros", 0, False),
])

# hipoteses
s = add_slide()
header(s, "Pergunta 4", "Hipóteses do projeto")
bullets(s, Inches(0.55), Inches(1.8), Inches(12.2), Inches(5.0), [
    ("H1 — Contratos de leasing com maior duração tendem a apresentar custos de manutenção mais elevados", 0, False),
    ("H2 — Carretas com maior tempo de utilização (idade) tendem a demandar maiores gastos com manutenção", 0, False),
    ("H3 — O aumento da quilometragem percorrida está associado ao aumento dos custos", 0, False),
    ("H4 — O histórico de manutenções anteriores é relevante para prever custos futuros", 0, False),
    ("H5 — Variáveis operacionais e características dos contratos influenciam significativamente os custos de manutenção", 0, False),
    ("As hipóteses são avaliadas na EDA (correlações e segmentações) e na modelagem (importância de variáveis)", 0, True),
])

# artigos
artigos = [
    ("Katreddi, Thiruvengadam, Thompson, Schmid & Padmanaban (2023)",
     "Machine learning models for maintenance cost estimation in delivery trucks using diesel and natural gas fuels",
     ["Prevê custo de manutenção por milha em caminhões de entrega; custos de manutenção são parcela relevante do TCO",
      "Compara Random Forest, XGBoost, ANN e ensemble Super Learner (variáveis: quilometragem, combustível, região, uso)",
      "Super Learner com melhor desempenho (R² = 97,28%; MAE = US$ 0,0068/milha) — ensembles capturam relações não lineares",
      "Relação com o projeto: variáveis operacionais semelhantes (km acumulado, região, montadora/modelo); reforça a escolha de modelos de árvore/ensemble; métricas devem ser lidas no contexto da nossa amostra"]),
    ("Katreddi, Thiruvengadam, Thompson & Schmid (2023)",
     "Mixed Effects Random Forest Model for Maintenance Cost Estimation in Heavy-Duty Vehicles Using Diesel and Alternative Fuels",
     ["Estima custos em veículos pesados com múltiplos combustíveis; tipo de veículo, região e operação influenciam custos",
      "Mixed Effects Random Forest: combina Random Forest com efeitos mistos por grupo de veículos",
      "Melhor generalização que Random Forest convencional em frotas heterogêneas",
      "Relação com o projeto: justifica incluir variáveis de agrupamento (montadora, modelo, classe, tipo de contrato) como features; MERF registrado como extensão futura"]),
    ("Sun, Guo, Sun, Yang & Hao (2024)",
     "Maintenance cost prediction for the vehicle based on maintenance data",
     ["Prevê custos futuros usando registros históricos de manutenção e falhas, com foco em garantia estendida",
      "Engenharia de confiabilidade: modelo Mixed Weibull + estimativa iterativa de custos (sem ML tradicional)",
      "Histórico de manutenção permite estimar custos futuros de forma consistente",
      "Relação com o projeto: sustenta o uso do histórico de OS como insumo preditivo, mesmo sem dados estruturados de falha por componente"]),
    ("Adekitan, Adetokun & Okokpujie (2018)",
     "A data-based investigation of vehicle maintenance cost components using ANN",
     ["Investiga fatores de custo de manutenção em veículos corporativos",
      "Rede Neural Artificial com quilometragem, consumo, frequência de falhas e histórico de uso",
      "Correlação R = 0,766 entre previsto e observado — sinal preditivo presente, desempenho modesto",
      "Relação com o projeto: mostra que variáveis operacionais carregam sinal, mas favorece priorizar métodos de árvore, que performaram melhor nos artigos mais recentes"]),
]
for i, (autores, titulo, pontos) in enumerate(artigos, 1):
    s = add_slide()
    header(s, f"Pergunta 5 · Artigo {i}/4", autores)
    tb(s, Inches(0.55), Inches(1.62), Inches(12.2), Inches(0.75), titulo, size=16, color=BLUE, bold=True)
    bullets(s, Inches(0.55), Inches(2.5), Inches(12.2), Inches(4.4),
            [(p, 0, j == 3) for j, p in enumerate(pontos)], size=16, gap=10)

# ============================================================ BLOCO 2
divider("Bloco 2", "Base de dados, variável-alvo e feature engineering",
        "Perguntas 6–7 da agenda da disciplina")

# base de dados
s = add_slide()
header(s, "Pergunta 6", "Base de dados — modelo estrela (1 dimensão + 6 fatos)")
tabela = [
    ("dim_carretas", "uma carreta", "atributos do ativo: montadora, modelo, ano, eixos, comprimento, reefer, classe"),
    ("fato_readings", "uma leitura de odômetro", "quilometragem acumulada por data"),
    ("fato_wo", "uma ordem de serviço", "cabeçalho da OS + totais internos"),
    ("fato_wo_labour", "uma linha de mão de obra", "custo interno de mão de obra + sistema VMRS"),
    ("fato_wo_parts", "uma linha de peça", "custo interno de peças + flag de garantia"),
    ("fato_contratos", "uma carreta-contrato", "leasing/rental: tipo, franquia de km, vigência"),
    ("fato_gps", "uma posição por dia", "telemetria: lat/long da carreta"),
]
yy = 1.85
tb(s, Inches(0.55), Inches(yy), Inches(2.9), Inches(0.35), "TABELA", size=13, color=ORANGE, bold=True)
tb(s, Inches(3.5), Inches(yy), Inches(3.0), Inches(0.35), "GRÃO (1 LINHA =)", size=13, color=ORANGE, bold=True)
tb(s, Inches(6.6), Inches(yy), Inches(6.2), Inches(0.35), "PAPEL", size=13, color=ORANGE, bold=True)
yy += 0.45
for nome, grao, papel in tabela:
    tb(s, Inches(0.55), Inches(yy), Inches(2.9), Inches(0.5), nome, size=14, color=NAVY, bold=True)
    tb(s, Inches(3.5), Inches(yy), Inches(3.0), Inches(0.5), grao, size=13, color=DARK)
    tb(s, Inches(6.6), Inches(yy), Inches(6.2), Inches(0.5), papel, size=13, color=GREY)
    yy += 0.52
tb(s, Inches(0.55), Inches(6.3), Inches(12.2), Inches(1.0),
   "Em palavras simples: são 7 \"planilhas\" ligadas pelo número da carreta (id_carreta) — uma com o cadastro de cada carreta e seis com os eventos da vida dela (leituras de km, ordens de serviço, custos de mão de obra e peças, contratos e posições de GPS).\n"
   "Detalhe técnico: mão de obra e peças também se ligam pela OS (id_os); contratos, pelo período de vigência. Extração: data/extract_custo_interno_km.sql, 2020–2025, frota própria.",
   size=11.5, color=GREY)

# custo interno
s = add_slide()
header(s, "Pergunta 6", "Definição do custo interno (o que entra no Y)")
bullets(s, Inches(0.55), Inches(1.8), Inches(12.2), Inches(3.4), [
    ("Em palavras simples: é o dinheiro que sai do caixa da empresa para manter a carreta rodando — a parte do conserto que NÃO é cobrada do cliente", 0, True),
    ("Tecnicamente: linhas de OS com charge_flag = 'I' (interno)", 0, False),
    ("Somente OS aprovadas, concluídas e não canceladas", 0, False),
    ("Mão de obra: se terceirizada (sublet_flag='Y') usa total_sublet; senão cost_hours × hourly_cost", 0, False),
    ("Peças: se terceirizada usa total_sublet; senão nvl(item_average_cost, item_cost) × actual_qty", 0, False),
    ("Inclui manutenção preventiva E corretiva — 'interno' não é sinônimo de 'preventivo'", 0, True),
])
card = s.shapes.add_shape(1, Inches(0.55), Inches(5.1), Inches(12.2), Inches(1.6))
card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.color.rgb = ORANGE; card.shadow.inherit = False
tb(s, Inches(0.9), Inches(5.3), Inches(11.6), Inches(1.2),
   f"Reconciliação: CAD {fmt(stats['custo_total_nominal']/1e6, 1)} milhões de custo interno nominal na janela 2020–2025 "
   f"(CAD {fmt(stats['custo_total_deflacionado']/1e6, 1)} mi em valores de dez/2025), distribuído em 238.818 ordens de serviço.",
   size=15, color=NAVY, bold=True)

# VMRS
s = add_slide()
header(s, "Contexto técnico", "VMRS — padrão de codificação de manutenção")
bullets(s, Inches(0.55), Inches(1.8), Inches(6.4), Inches(4.8), [
    ("Em palavras simples: é o \"CID da oficina\" — em vez de descrever o conserto em texto livre, cada reparo recebe um código padronizado que diz qual sistema da carreta foi mexido", 0, True),
    ("Nome oficial: Vehicle Maintenance Reporting System (ATA)", 0, False),
    ("Usamos o nível de sistema — suficiente para saber qual sistema gera o custo", 0, False),
    ("Permite agregar custo por sistema de forma comparável entre todas as carretas", 0, False),
    ("Importante: 'PM' (manutenção preventiva) é apenas UMA das categorias VMRS — ao lado de freios, pneus, reefer, elétrica etc.", 0, True),
    ("A variável-alvo do projeto soma o custo interno de TODAS as categorias (preventivas e corretivas) — o VMRS entra como dimensão de análise, não como filtro do Y", 0, True),
])
vmrs_ex = [("PM", "Preventive Maintenance"), ("04", "Brakes"), ("09", "Tires and Accessories"),
           ("10", "Reefer"), ("03", "Lights and Wiring"), ("08", "Exterior Body"),
           ("05", "Landing Gear"), ("02", "Air Equipment")]
yy = 2.0
tb(s, Inches(7.3), Inches(1.8), Inches(5.4), Inches(0.35), "EXEMPLOS DE SISTEMAS", size=13, color=ORANGE, bold=True)
yy = 2.25
for cod, nome in vmrs_ex:
    tb(s, Inches(7.3), Inches(yy), Inches(1.0), Inches(0.4), cod, size=14, color=NAVY, bold=True)
    tb(s, Inches(8.3), Inches(yy), Inches(4.4), Inches(0.4), nome, size=13, color=DARK)
    yy += 0.45

# programas PM
s = add_slide()
header(s, "Contexto técnico", "Programas de manutenção preventiva da frota")
progs = [
    ("PM · Inspeção geral", "Estado mecânico e de segurança: chassi, landing gear, pneus e torque, freios, luzes/ABS, portas e vedações; anotação de defeitos para OS"),
    ("Inspeção de Segurança (MTO)", "Inspeção regulatória anual: freios completos, sistema de ar, suspensão, rodas e cubos, chassi e coupling, elétrica — emite o certificado e selo"),
    ("Reefer PM · Dry Service", "Serviço 'seco' da unidade de refrigeração (Thermo King/Carrier): limpeza de condensador/evaporador, correias, teste de ciclo, sensores e datalogger"),
    ("Reefer PM · Wet Service", "Dry Service + trocas de fluidos: óleo e filtros da unidade, filtros de combustível/ar, coolant, correia; reset do horímetro"),
]
yy = 1.9
for nome, desc_ in progs:
    card = s.shapes.add_shape(1, Inches(0.55), Inches(yy), Inches(12.2), Inches(1.12))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.color.rgb = RGBColor(0xDD, 0xD5, 0xC8); card.shadow.inherit = False
    tb(s, Inches(0.85), Inches(yy + 0.08), Inches(3.6), Inches(0.9), nome, size=15, color=NAVY, bold=True)
    tb(s, Inches(4.5), Inches(yy + 0.08), Inches(8.0), Inches(1.0), desc_, size=12.5, color=DARK)
    yy += 1.28
tb(s, Inches(0.55), Inches(7.0), Inches(12.2), Inches(0.4),
   "Além dos PMs programados, manutenções corretivas não programadas (luzes, freios, pneus...) também geram custo interno — e são as que mais variam entre carretas.",
   size=12, color=GREY)

# variavel Y
s = add_slide()
header(s, "Variável-alvo (Y)", "custo_manutencao_interno_por_km")
card = s.shapes.add_shape(1, Inches(0.55), Inches(1.85), Inches(12.2), Inches(1.15))
card.fill.solid(); card.fill.fore_color.rgb = NAVY; card.line.fill.background(); card.shadow.inherit = False
tb(s, Inches(0.95), Inches(2.12), Inches(11.4), Inches(0.7),
   "Y = custo interno total do mês (CAD, deflacionado)  ÷  km rodado no mês   ·   grão: carreta × mês",
   size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
bullets(s, Inches(0.55), Inches(3.25), Inches(12.2), Inches(3.8), [
    ("Em palavras simples: para cada carreta, em cada mês, dividimos o que a empresa gastou com ela pelo quanto ela rodou. O resultado responde: \"quanto custa manter esta carreta, por km rodado?\"", 0, True),
    ("Numerador: mão de obra + peças internas das OS do mês, de TODAS as categorias VMRS — preventiva (PM) é só uma delas; corretivas (freios, pneus, reefer...) também entram", 0, False),
    ("Denominador: km rodado no mês, calculado a partir das leituras de odômetro (resets e saltos anômalos tratados); só calculamos quando a carreta rodou ≥ 500 km no mês", 0, False),
    (f"População com Y válido: {int(stats['obs_y_valido']):,} observações carreta × mês".replace(",", "."), 0, False),
    (f"Em {stats['share_y_zero']*100:.1f}% dos meses a carreta nem passou pela oficina (custo zero) — e isso é informação legítima para orçamento, então esses meses foram mantidos", 0, True),
    (f"Média: CAD {fmt(stats['y_media'])} /km · mediana dos meses com custo: CAD {fmt(stats['y_mediana_positivos'])} /km", 0, False),
])

# deflacao
s = add_slide()
header(s, "Correção metodológica", "Custos em CAD ⇒ deflação pelo CPI do Canadá")
bullets(s, Inches(0.55), Inches(1.75), Inches(5.9), Inches(5.2), [
    ("Em palavras simples: 1 dólar de 2020 comprava mais coisas que 1 dólar de 2025. Para comparar custos de anos diferentes de forma justa, convertemos tudo para o \"dólar de dez/2025\"", 0, True),
    ("Os custos da operação são em dólares canadenses (CAD) — por isso usamos o índice de inflação do Canadá (CPI, Statistics Canada), corrigindo a versão anterior que usava IPCA (Brasil)", 0, False),
    ("A inflação canadense acumulada 2020→2025 foi ≈ 20% — sem a correção, um modelo confundiria \"custos subindo\" com \"dinheiro valendo menos\"", 0, True),
    ("Mecânica: cada custo é multiplicado pelo índice de dez/2025 dividido pelo índice do mês em que ocorreu", 0, False),
])
pic_fit(s, FIG / "04_nominal_vs_deflacionado.png", Inches(6.6), Inches(1.9), Inches(6.3), Inches(4.7))

# normalizado vs desnormalizado
s = add_slide()
header(s, "Modelagem de dados", "Do modelo estrela à base de modelagem (carreta × mês)")
bullets(s, Inches(0.55), Inches(1.8), Inches(6.2), Inches(5.0), [
    ("Em palavras simples: juntamos as 7 planilhas numa única tabela grande, em que cada linha responde \"como estava a carreta X no mês Y?\" — o formato que a análise e os modelos conseguem consumir", 0, True),
    ("Normalizado (origem): 7 tabelas relacionadas, sem repetição — formato da extração do ERP", 0, False),
    ("Desnormalizado (análise): uma tabela larga, uma linha = uma carreta em um mês, com todas as variáveis juntas", 0, False),
    ("Pipeline: join por id_carreta (e período, para contratos) + feature engineering + agregação mensal", 0, False),
    (f"Resultado: base_mensal_carreta com {int(stats['linhas_base_total']):,} linhas (grade completa carreta × mês 2020–2025)".replace(",", "."), 0, False),
    ("Redundância aceita (ex.: montadora repete a cada mês) em troca de uma base pronta para correlação e treino de modelos", 0, False),
])
steps = ["7 tabelas normalizadas\n(extração · modelo estrela)",
         "join + feature engineering\n+ agregação mensal",
         "1 tabela desnormalizada\n(base carreta × mês · EDA & ML)"]
yy = 2.2
for i, txt in enumerate(steps):
    card = s.shapes.add_shape(1, Inches(7.2), Inches(yy), Inches(5.4), Inches(1.05))
    card.fill.solid(); card.fill.fore_color.rgb = NAVY if i != 1 else ORANGE
    card.line.fill.background(); card.shadow.inherit = False
    tb(s, Inches(7.45), Inches(yy + 0.14), Inches(4.9), Inches(0.85), txt, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    yy += 1.45

# feature engineering
s = add_slide()
header(s, "Pergunta 7", "Feature engineering — variáveis derivadas")
fe = [
    ("idade_carreta", "anos desde a entrada em serviço"),
    ("km_rodado_mes / km_por_mes", "intensidade de uso mensal (prorrateio do odômetro)"),
    ("km_acumulado (defasado)", "quilometragem acumulada até o mês anterior"),
    ("custo_acum_manutencao", "custo interno acumulado até o mês anterior"),
    ("n_os_acum / n_os_preventivas_acum", "nº de OS acumuladas (totais e preventivas)"),
    ("custo_medio_movel_3m", "média móvel de custo dos 3 meses anteriores"),
    ("intervalo_medio_os", "dias médios entre OS da carreta"),
    ("meses_desde_ultima_os", "recência da última OS (apenas OS anteriores ao mês)"),
    ("duracao_contrato_meses", "duração do contrato vigente"),
    ("idade_contrato_meses_no_mes", "tempo decorrido de contrato no mês"),
    ("regiao_operacao", "região aproximada por província/local da OS"),
    ("custos deflacionados (CPI Canadá)", "todos os custos a valor de dez/2025"),
]
yy = 1.85
col_x = [Inches(0.55), Inches(6.85)]
for i, (nome, desc_) in enumerate(fe):
    x = col_x[i % 2]
    if i % 2 == 0 and i > 0:
        yy += 0.78
    tb(s, x, Inches(yy), Inches(5.9), Inches(0.4), nome, size=14, color=NAVY, bold=True)
    tb(s, x, Inches(yy + 0.32), Inches(5.9), Inches(0.4), desc_, size=12, color=GREY)
tb(s, Inches(0.55), Inches(6.7), Inches(12.2), Inches(0.75),
   "Em palavras simples: criamos colunas novas que resumem a \"vida\" da carreta até aquele mês — quanto já rodou, quanto já custou, quando foi a última visita à oficina. Regra de ouro (anti-vazamento): cada linha só usa informação que já existia ANTES do mês que queremos prever — como na vida real. prop_pecas_garantia foi descartada: apenas 0,03% dos meses têm peça em garantia.",
   size=11.5, color=GREY)

# ---- dicionario de variaveis do modelo (naturais vs FE) ----
def slide_dicionario(titulo, linhas, col_w=(3.55, 1.15, 7.5)):
    s = add_slide()
    header(s, "Pergunta 6–7", titulo)
    xs = [Inches(0.55), Inches(0.55 + col_w[0]), Inches(0.55 + col_w[0] + col_w[1])]
    tb(s, xs[0], Inches(1.72), Inches(col_w[0]), Inches(0.3), "VARIÁVEL", size=12, color=ORANGE, bold=True)
    tb(s, xs[1], Inches(1.72), Inches(col_w[1]), Inches(0.3), "TIPO", size=12, color=ORANGE, bold=True)
    tb(s, xs[2], Inches(1.72), Inches(col_w[2]), Inches(0.3), "SIGNIFICADO · COMPOSIÇÃO (QUANDO FE)", size=12, color=ORANGE, bold=True)
    yy = 2.08
    for nome, tipo, desc_ in linhas:
        cor_tipo = BLUE if tipo == "Natural" else ORANGE
        tb(s, xs[0], Inches(yy), Inches(col_w[0]), Inches(0.55), nome, size=12.5, color=NAVY, bold=True)
        tb(s, xs[1], Inches(yy), Inches(col_w[1]), Inches(0.55), tipo, size=11.5, color=cor_tipo, bold=True)
        tb(s, xs[2], Inches(yy), Inches(col_w[2]), Inches(0.6), desc_, size=11, color=DARK)
        yy += 0.555
    return s

slide_dicionario("Dicionário das variáveis do modelo — numéricas (1/2)", [
    ("km_rodado_mes", "FE", "SUBTRAÇÃO + RATEIO: leitura atual do odômetro − leitura anterior; se o intervalo cruza meses, o delta é dividido proporcionalmente aos dias de cada mês (resets/saltos removidos antes)"),
    ("km_acumulado", "FE", "ÚLTIMA LEITURA DEFASADA: maior leitura de odômetro disponível até o fim do MÊS ANTERIOR (nunca do mês previsto — anti-vazamento)"),
    ("km_rodado_acum", "FE", "SOMA ACUMULADA: Σ km_rodado_mes de todos os meses anteriores da carreta (exclui o mês corrente)"),
    ("km_por_mes", "FE", "DIVISÃO: km_acumulado ÷ nº de meses desde data_entrada_servico"),
    ("idade_carreta", "FE", "SUBTRAÇÃO: (mês de referência − data_entrada_servico) ÷ 365,25 → idade em anos"),
    ("custo_acum_manutencao", "FE", "SOMA ACUMULADA DEFASADA: Σ custo_total_mes dos meses anteriores (cumsum − valor do mês corrente)"),
    ("custo_preventivo_acum", "FE", "SOMA ACUMULADA DEFASADA: idem, somando apenas custo_preventivo_total_mes (linhas VMRS PM)"),
    ("n_os_acum", "FE", "CONTAGEM ACUMULADA: Σ nº de OS dos meses anteriores da carreta"),
])

slide_dicionario("Dicionário das variáveis do modelo — numéricas (2/2)", [
    ("n_os_preventivas_acum", "FE", "CONTAGEM ACUMULADA: Σ nº de OS preventivas dos meses anteriores"),
    ("custo_medio_movel_3m", "FE", "MÉDIA MÓVEL: (custo do mês−1 + mês−2 + mês−3) ÷ 3 — janela deslocada 1 mês para trás (shift), sem incluir o mês previsto"),
    ("custo_preventivo_medio_movel_3m", "FE", "MÉDIA MÓVEL: mesma fórmula, sobre o custo preventivo dos 3 meses anteriores"),
    ("intervalo_medio_os", "FE", "MÉDIA DE DIFERENÇAS: média dos dias entre OS consecutivas da carreta, usando só OS anteriores ao mês previsto"),
    ("meses_desde_ultima_os", "FE", "SUBTRAÇÃO: mês de referência − mês da última OS ANTERIOR (nunca conta OS do próprio mês — anti-vazamento)"),
    ("franquia_km_mensal", "Natural", "campo direto do contrato vigente (fato_contratos) — sem transformação"),
    ("duracao_contrato_meses", "FE", "SUBTRAÇÃO: (data_fim − data_inicio) do contrato ÷ 30,44 → meses (contratos abertos: censurados em 2025-12)"),
    ("idade_contrato_meses_no_mes", "FE", "SUBTRAÇÃO: (mês de referência − data_inicio do contrato) ÷ 30,44 → meses decorridos"),
])

slide_dicionario("Dicionário das variáveis do modelo — categóricas", [
    ("cod_montadora", "Natural", "fabricante da carreta (dim_carretas) — one-hot com agrupamento de categorias raras"),
    ("cod_modelo", "Natural", "modelo da carreta (dim_carretas)"),
    ("flag_refrigerado", "Natural", "carreta refrigerada (reefer) vs seca (dim_carretas)"),
    ("tipo_contrato", "Natural", "RENTAL vs LEASE do contrato vigente (fato_contratos)"),
    ("cod_grupo_manutencao", "Natural", "programa/grupo de manutenção do ativo (dim_carretas)"),
    ("regiao_operacao", "FE", "MAPEAMENTO + PROPAGAÇÃO: província/local da OS → região; carrega-se a última região conhecida para os meses seguintes; na modelagem usa-se a região do mês ANTERIOR (defasada)"),
])
s = prs.slides[-1]
tb(s, Inches(0.55), Inches(5.6), Inches(12.2), Inches(1.3),
   "Notas: variáveis 'FE' (feature engineering) são derivadas das tabelas-fonte; todas as históricas usam apenas informação disponível até o mês anterior. "
   "Excluídas da modelagem: prop_pecas_garantia (apenas 0,03% dos meses com peça em garantia), tipo_manutencao (fixada em MAINT como população) e km_acumulado_fim_mes (auditoria).",
   size=12, color=GREY)

# ============================================================ EDA
divider("Bloco 2 · Pergunta 8", "Análise Exploratória de Dados (EDA)",
        f"{int(stats['obs_y_valido']):,} observações carreta × mês com Y válido · protocolo variável a variável".replace(",", "."))

# protocolo
s = add_slide()
header(s, "Pergunta 8", "Protocolo da EDA")
bullets(s, Inches(0.55), Inches(1.8), Inches(12.2), Inches(5.0), [
    ("Para cada variável quantitativa (X e Y):", 0, True),
    ("histograma (distribuição, assimetria) · boxplot (outliers, dispersão) · estatísticas: N, média, DP, min, Q1, mediana, Q3, max", 1, False),
    ("Para cada variável qualitativa:", 0, True),
    ("boxplot de Y segmentado por categoria · tabela de frequência · estatísticas de Y por categoria", 1, False),
    ("Análises complementares:", 0, True),
    ("correlação de Pearson e Spearman de cada X com Y · força de associação (eta) para categóricas · matriz de correlação entre X · multicolinearidade (VIF) · evolução temporal", 1, False),
    ("Leitura em duas partes (zero-inflação):", 0, True),
    ("ocorrência de custo no mês (67% dos meses = zero) e magnitude condicional nos meses com custo > 0", 1, False),
])

# como ler os graficos
s = add_slide()
header(s, "Pergunta 8", "Como ler os gráficos desta seção")
bullets(s, Inches(0.55), Inches(1.75), Inches(12.2), Inches(5.3), [
    ("Histograma — mostra a distribuição: onde os valores se concentram e se há cauda (assimetria). Cauda longa à direita = poucos casos muito caros puxando a média para cima.", 0, False),
    ("Boxplot — resume a distribuição: a caixa vai do 1º ao 3º quartil (50% central dos dados), a linha vermelha é a mediana e os pontos isolados são outliers. Caixa 'achatada' perto do zero = maioria dos valores baixos.", 0, False),
    ("Correlação de Spearman (ρ) — mede se 'quanto mais X, mais Y' (de −1 a +1). Usamos Spearman porque é robusta a outliers; Pearson é reportada como referência. Regra de leitura: |ρ| < 0,1 fraca · 0,1–0,3 modesta · > 0,3 relevante.", 0, False),
    ("Eta (η) para categóricas — mede o quanto as categorias separam o Y (0 = categorias iguais; 1 = categoria determina o custo).", 0, False),
    ("Gráficos cortados no percentil 99 (p99) — 1% dos meses mais caros é omitido DO GRÁFICO para dar legibilidade à massa central; esses outliers são tratados na modelagem, não ignorados.", 0, False),
    ("Em cada slide: figura à esquerda, estatísticas à direita, e a leitura (o que o gráfico diz) na faixa azul inferior.", 0, True),
])

# --- Y ---
s = add_slide()
header(s, "EDA · variável-alvo", "Distribuição do custo interno por km (deflacionado)")
pic_fit(s, FIG_EDA / f"quant_{TARGET}.png", Inches(0.4), Inches(1.7), Inches(8.6), Inches(4.3))
drow = desc[desc["variavel"] == TARGET].iloc[0]
stattxt = (f"N = {int(drow['N']):,}".replace(",", ".") + "\n"
           f"média = {fmt(drow['media'])}\nDP = {fmt(drow['desvio_padrao'])}\n"
           f"mín = {fmt(drow['min'])}\nQ1 = {fmt(drow['Q1'])}\nmediana = {fmt(drow['mediana'])}\n"
           f"Q3 = {fmt(drow['Q3'])}\nmáx = {fmt(drow['max'], 1)}\nassimetria = {fmt(drow['assimetria'], 1)}")
tb(s, Inches(9.3), Inches(1.9), Inches(3.6), Inches(0.4), "ESTATÍSTICAS (CAD/km)", size=13, color=ORANGE, bold=True)
tb(s, Inches(9.3), Inches(2.3), Inches(3.6), Inches(3.4), stattxt, size=14, color=DARK, line_spacing=1.15)
tb(s, Inches(0.55), Inches(6.2), Inches(12.2), Inches(1.0),
   f"Distribuição fortemente assimétrica à direita e zero-inflada: {stats['share_y_zero']*100:.0f}% dos meses sem custo interno. "
   f"Nos meses com custo, a mediana é CAD {fmt(stats['y_mediana_positivos'])}/km e a média CAD {fmt(stats['y_media_positivos'])}/km — cauda longa de meses caros puxa a média. "
   "Implicação: transformação log1p e modelos robustos; o problema combina ocorrência + magnitude.",
   size=13, color=NAVY, bold=True)

# evolucao Y
s = add_slide()
header(s, "EDA · variável-alvo", "Evolução anual do custo interno por km (valores reais)")
pic_fit(s, FIG_EDA / "evolucao_y_anual.png", Inches(0.4), Inches(1.7), Inches(8.6), Inches(4.4))
e0, e1 = evo_y.iloc[0], evo_y.iloc[-1]
tb(s, Inches(9.3), Inches(1.9), Inches(3.7), Inches(4.6),
   f"2020: CAD {fmt(e0['media_y'])}/km\n2025: CAD {fmt(e1['media_y'])}/km\n\n"
   f"Crescimento real de {((e1['media_y']/e0['media_y'])-1)*100:.0f}% na média — mesmo após remover a inflação (CPI Canadá).\n\n"
   f"Share de meses com custo zero caiu de {e0['share_zero']*100:.0f}% para {e1['share_zero']*100:.0f}%: a frota está gerando OS com mais frequência.",
   size=14, color=DARK, line_spacing=1.15)
tb(s, Inches(0.55), Inches(6.3), Inches(12.2), Inches(0.8),
   "Leitura: há tendência real de encarecimento por km — consistente com envelhecimento/maior utilização da frota — que o modelo deve capturar pelas variáveis operacionais, não por um simples efeito 'ano'.",
   size=13, color=NAVY, bold=True)

# --- quantitativas: 1 slide por variavel ---
quant_coment = {
    "ano_modelo": "O que o gráfico mostra: a frota se concentra em modelos 2014–2022, com picos nos anos de maior renovação. A correlação com o custo por km é quase nula (ρ = +0,04): saber só o ano do modelo não diz quanto a carreta custa. O efeito 'envelhecimento' existe, mas aparece nas variáveis de histórico acumulado — por isso o ano entra no modelo apenas como apoio.",
    "eixos": "O que o gráfico mostra: quase toda a frota tem 2 eixos — a barra única domina o histograma e o boxplot é uma linha. Quando uma variável quase não varia, ela não consegue separar carretas caras de baratas (ρ ≈ 0). Conclusão: pouco útil para o modelo, mantida apenas por completude.",
    "comprimento": "O que o gráfico mostra: o padrão de 53 pés domina a frota — a distribuição é quase toda um único valor. Com tão pouca variação, a associação com o custo é fraca (ρ = −0,07). Leitura de negócio: o comprimento não diferencia custo porque praticamente todas as carretas são do mesmo porte.",
    "idade_carreta": "O que o gráfico mostra: idade média de ≈ 6 anos, com cauda de carretas antigas (até 28 anos). A correlação direta com o custo mensal por km é modesta (ρ = +0,04) — surpreendente à primeira vista. Explicação: a idade age INDIRETAMENTE, via km acumulado e histórico de OS (carretas velhas rodaram mais e abriram mais OS); essas variáveis absorvem o efeito. H2 tem suporte apenas parcial.",
    "km_rodado_mes": "O que o gráfico mostra: metade das carretas roda até ≈ 1.900 km/mês, mas há cauda de alta utilização (> 10 mil km/mês). Cuidado na leitura: esta variável é o DENOMINADOR do Y (custo ÷ km) — rodar mais dilui o custo por km mecanicamente. Por isso a correlação com custo/km não deve ser lida como efeito causal; a hipótese H3 foi testada também contra o custo absoluto.",
    "km_acumulado": "O que o gráfico mostra: forte assimetria — metade da frota tem até ≈ 86 mil km acumulados, mas existem carretas com milhões de km. Correlação positiva modesta com o custo (ρ = +0,16): quanto mais rodada a carreta, mais cara por km. É o sinal de desgaste estrutural que sustenta parcialmente a H3.",
    "km_por_mes": "O que o gráfico mostra: intensidade média de uso desde a entrada em serviço — mediana de ≈ 900 km/mês, com carretas de uso intenso na cauda. ρ = +0,12: carretas que historicamente rodam mais tendem a custar mais por km. Complementa km_rodado_mes com uma visão estrutural (perfil da carreta) em vez de conjuntural (o mês específico).",
    "franquia_km_mensal": "O que o gráfico mostra: a grande maioria dos contratos não tem franquia de km (valor zero) — a distribuição é quase toda um único valor. Sem variação, não há poder preditivo (ρ ≈ 0). Leitura de negócio: a franquia contratual, quando existe, é exceção — não é um driver de custo nesta frota.",
    "duracao_contrato_meses": "O que o gráfico mostra: contratos tipicamente de 60 meses (5 anos), o padrão do mercado de leasing. A correlação com o custo por km é praticamente nula (ρ = +0,02). Conclusão importante: a H1 (contratos mais longos ⇒ mais custo) NÃO encontra suporte — a duração do contrato, por si só, não explica o custo de manutenção.",
    "idade_contrato_meses_no_mes": "O que o gráfico mostra: tempo decorrido de contrato no mês da observação — distribuição espalhada, refletindo contratos em fases diferentes. Associação fraca com o custo (ρ = +0,05): estar no início ou no fim do contrato pouco muda o custo por km. Reforça a leitura da H1: o efeito contratual é pequeno.",
    "custo_acum_manutencao": "O que o gráfico mostra: forte assimetria — metade das carretas acumulou até CAD ≈ 2 mil de custo interno, mas a cauda chega a milhões. ρ = +0,20, o 2º preditor mais forte: carreta que já custou tende a continuar custando. É a evidência central da H4 — o passado de manutenção prevê o futuro.",
    "custo_preventivo_acum": "O que o gráfico mostra: custo preventivo acumulado com o mesmo padrão de cauda longa. ρ = +0,17: quem consumiu mais preventiva também custa mais por km no total — em parte porque carretas mais usadas fazem mais PMs. Reforça o bloco de histórico (H4).",
    "n_os_acum": "O que o gráfico mostra: metade das carretas acumulou até 9 OS; a cauda passa de centenas. ρ = +0,22 — o preditor individual MAIS FORTE do custo por km. Leitura de negócio: a frequência histórica de manutenção é o melhor termômetro isolado de quanto a carreta custará — melhor que idade, modelo ou contrato. Núcleo da H4.",
    "n_os_preventivas_acum": "O que o gráfico mostra: OS preventivas acumuladas — mediana de 5, máximo de 38 (os PMs são programados, então o teto é limitado). ρ = +0,19, quase tão forte quanto n_os_acum — mas as duas andam juntas (carreta com mais OS tem mais PMs). Essa redundância é tratada na análise de colinearidade.",
    "custo_medio_movel_3m": "O que o gráfico mostra: o custo médio dos 3 meses anteriores é zero ou baixo para a maioria (reflexo dos meses sem OS), com cauda de carretas em 'sequência cara'. ρ = +0,08: há persistência de curto prazo — um trimestre caro tende a ser seguido de custo acima da média. Sinal moderado, mas útil combinado às demais.",
    "custo_preventivo_medio_movel_3m": "O que o gráfico mostra: média móvel preventiva com distribuição semelhante, dominada por zeros. Associação fraca com o custo total (ρ = −0,02): fazer preventiva no trimestre anterior não aumenta nem reduz visivelmente o custo por km do mês seguinte — no grão mensal, o retorno da preventiva não é imediato.",
    "intervalo_medio_os": "O que o gráfico mostra: em média, as carretas abrem OS a cada ≈ 100 dias (mediana 81). A correlação é NEGATIVA (ρ = −0,19): intervalos curtos = manutenção frequente = custo/km alto. É o espelho de n_os_acum e o 3º sinal mais forte — carreta que volta cedo à oficina é carreta cara.",
    "meses_desde_ultima_os": "O que o gráfico mostra: metade das observações está a até 2 meses da última OS. Correlação fraca e negativa (ρ = −0,06): quem fez OS há pouco tende a custar um pouco mais no mês seguinte (problemas vêm em sequência). Importante: a variável usa apenas OS ANTERIORES ao mês previsto — a versão antiga vazava informação do próprio mês e inflava as métricas.",
}
for _, row in desc[desc["tipo"] == "Quantitativa"].iterrows():
    col = row["variavel"]
    fig_path = FIG_EDA / f"quant_{col}.png"
    if not fig_path.exists():
        continue
    s = add_slide()
    header(s, "EDA · quantitativa", col)
    pic_fit(s, fig_path, Inches(0.4), Inches(1.6), Inches(8.6), Inches(3.75))
    crow = corr[corr["variavel"] == col]
    ctxt = ""
    if not crow.empty:
        ctxt = f"\n\nPearson c/ Y = {fmt(crow.iloc[0]['pearson'])}\nSpearman c/ Y = {fmt(crow.iloc[0]['spearman'])}"
    stattxt = (f"N = {int(row['N']):,}".replace(",", ".") + "\n"
               f"média = {fmt(row['media'], 2)}\nDP = {fmt(row['desvio_padrao'], 2)}\n"
               f"mín = {fmt(row['min'], 2)}\nQ1 = {fmt(row['Q1'], 2)}\nmediana = {fmt(row['mediana'], 2)}\n"
               f"Q3 = {fmt(row['Q3'], 2)}\nmáx = {fmt(row['max'], 0)}" + ctxt)
    tb(s, Inches(9.3), Inches(1.7), Inches(3.6), Inches(0.4), "ESTATÍSTICAS", size=13, color=ORANGE, bold=True)
    tb(s, Inches(9.3), Inches(2.1), Inches(3.6), Inches(3.6), stattxt, size=12.5, color=DARK, line_spacing=1.1)
    card = s.shapes.add_shape(1, Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.75))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = NAVY; card.shadow.inherit = False
    tb(s, Inches(0.7), Inches(5.62), Inches(12.0), Inches(1.55),
       quant_coment.get(col, ""), size=12.5, color=NAVY, line_spacing=1.08)

# --- qualitativas: 1 slide por variavel ---
quali_coment = {
    "cod_montadora": "O que o gráfico mostra: cada caixa é a distribuição do custo por km das carretas de uma montadora, ordenadas da mediana mais alta para a mais baixa. Há diferença visível entre fabricantes, mas as caixas se sobrepõem bastante — o efeito isolado é fraco (η = 0,068). No modelo, a montadora entra via one-hot (categorias raras agrupadas) e contribui em interação com as demais variáveis.",
    "flag_refrigerado": "O que o gráfico mostra: a caixa das refrigeradas (reefer) está deslocada para cima em relação às secas — custam mais por km em toda a distribuição, não só nos extremos. Faz sentido operacional: a unidade de refrigeração adiciona um sistema inteiro (compressor, gás, filtros) que gera OS próprias. η = 0,063; a flag aparece entre os fatores importantes do modelo final.",
    "tipo_contrato": "O que o gráfico mostra: RENTAL e LEASE têm medianas próximas — o tipo de contrato desloca pouco o custo por km (η = 0,040). Leitura para a H5: características contratuais influenciam, mas bem menos do que o histórico operacional da carreta.",
    "tipo_manutencao": "O que o gráfico mostra: MAINT (manutenção inclusa no contrato) concentra 94,5% da população modelada — as caixas de NET/MIX refletem populações pequenas e regras de cobrança diferentes. Por isso a modelagem é feita SÓ com MAINT: comparar tipos misturados confundiria regra contratual com comportamento de custo.",
    "cod_grupo_manutencao": "O que o gráfico mostra: os programas de manutenção têm caixas muito parecidas — quase não separam custo por km (η = 0,016, o mais fraco entre as categóricas). O programa define QUANDO a preventiva acontece, mas não diferencia QUANTO a carreta custa no total.",
    "regiao_operacao": "O que o gráfico mostra: a região é a categórica mais forte (η = 0,084) — as medianas variam visivelmente entre províncias/regiões. Leitura de negócio: clima (inverno rigoroso), qualidade das estradas e custo de oficina local mudam o custo por km. Suporte parcial à H5 pelo lado operacional.",
    "cod_classe": "O que o gráfico mostra: as classes de ativo têm medianas distintas (η = 0,056) — classes mais complexas custam mais por km. Efeito real porém moderado; correlacionada com a montadora e a flag reefer (as classes agrupam tipos de carreta).",
}
for col in ["cod_montadora", "flag_refrigerado", "tipo_contrato", "tipo_manutencao",
            "cod_grupo_manutencao", "regiao_operacao", "cod_classe"]:
    fig_path = FIG_EDA / f"quali_{col}.png"
    if not fig_path.exists():
        continue
    s = add_slide()
    header(s, "EDA · qualitativa", col)
    pic_fit(s, fig_path, Inches(0.4), Inches(1.6), Inches(9.2), Inches(3.75))
    f_var = freq[freq["variavel"] == col].nlargest(6, "n")
    e_var = eta[eta["variavel"] == col]
    txt = "FREQUÊNCIA (TOP 6)\n" + "\n".join(
        f"{r['categoria'][:16]}: {r['pct']*100:.1f}%" for _, r in f_var.iterrows())
    if not e_var.empty:
        txt += f"\n\neta = {fmt(e_var.iloc[0]['eta'])}\ncategorias = {int(e_var.iloc[0]['n_categorias'])}"
    tb(s, Inches(9.8), Inches(1.7), Inches(3.2), Inches(3.8), txt, size=12, color=DARK, line_spacing=1.1)
    card = s.shapes.add_shape(1, Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.75))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = NAVY; card.shadow.inherit = False
    tb(s, Inches(0.7), Inches(5.62), Inches(12.0), Inches(1.55),
       quali_coment.get(col, ""), size=12.5, color=NAVY, line_spacing=1.08)

# ranking associacao
s = add_slide()
header(s, "EDA · seleção de variáveis", "Força de associação de cada X com o Y")
pic_fit(s, FIG_EDA / "ranking_associacao_y.png", Inches(0.4), Inches(1.6), Inches(8.4), Inches(5.2))
tb(s, Inches(9.0), Inches(1.9), Inches(4.0), Inches(4.8),
   "Nenhuma variável isolada é forte (máx ρ = 0,22).\n\n"
   "O sinal está no HISTÓRICO da carreta: nº de OS acumuladas, custo acumulado e intervalo entre OS dominam o ranking — suporte direto à H4.\n\n"
   "Idade, contrato e atributos físicos têm efeito fraco isoladamente (H1/H2 sem suporte direto no grão mensal).\n\n"
   "Implicação: o ganho preditivo depende de interações — favorece modelos de árvore/ensemble.",
   size=14, color=DARK, line_spacing=1.15)

# colinearidade — slide conceito + slide resultado
s = add_slide()
header(s, "EDA · multicolinearidade", "O que é multicolinearidade e por que verificamos")
bullets(s, Inches(0.55), Inches(1.8), Inches(12.2), Inches(5.2), [
    ("O problema: duas variáveis que 'dizem a mesma coisa' (ex.: idade da carreta e ano do modelo — uma é quase o espelho da outra)", 0, True),
    ("Por que importa: em modelos lineares, variáveis redundantes disputam o mesmo efeito e os coeficientes ficam instáveis — o modelo pode dizer que idade 'não importa' só porque o ano do modelo já capturou o efeito", 0, False),
    ("Como medimos:", 0, True),
    ("Matriz de correlação de Spearman entre as X — pares com |ρ| > 0,7 são candidatos a redundância", 1, False),
    ("VIF (Variance Inflation Factor) — mede o quanto cada X é 'explicável' pelas outras: VIF > 5 pede atenção, VIF > 10 indica redundância problemática", 1, False),
    ("O que fazemos quando há redundância: manter apenas uma variável do par (modelos lineares) ou usar modelos de árvore/ensemble, que são robustos a colinearidade — foi o caminho adotado", 0, True),
])

s = add_slide()
header(s, "EDA · multicolinearidade", "Resultado: quais variáveis se sobrepõem")
pic_fit(s, FIG_EDA / "matriz_spearman.png", Inches(0.3), Inches(1.55), Inches(7.4), Inches(4.0))
vif_top = vif.head(8)
txt = "VIF (TOP 8)\n" + "\n".join(
    f"{r['variavel'][:28]}: {r['vif']:.1f}" for _, r in vif_top.iterrows())
tb(s, Inches(8.0), Inches(1.7), Inches(4.9), Inches(3.6), txt, size=12.5, color=DARK, line_spacing=1.12)
card = s.shapes.add_shape(1, Inches(0.4), Inches(5.62), Inches(12.5), Inches(1.65))
card.fill.solid(); card.fill.fore_color.rgb = LIGHT
card.line.color.rgb = NAVY; card.shadow.inherit = False
tb(s, Inches(0.7), Inches(5.74), Inches(12.0), Inches(1.45),
   "Como ler: na matriz, laranja = correlação positiva, azul = negativa; quanto mais forte a cor, maior a sobreposição. "
   "Pares redundantes encontrados: nº de OS acumuladas ↔ nº de OS preventivas acumuladas, custo acumulado ↔ custo preventivo acumulado, idade ↔ ano do modelo. "
   "Decisão: como o modelo recomendado é de árvore (Random Forest), robusto a colinearidade, os pares foram mantidos; a restrição vale para os modelos lineares de comparação.",
   size=12.5, color=NAVY, line_spacing=1.08)

# outliers por variavel
try:
    outl = pd.read_csv(TABLES / "03d_diagnostico_outliers.csv")
except FileNotFoundError:
    outl = None
if outl is not None:
    s = add_slide()
    header(s, "EDA · outliers", "Diagnóstico de outliers por variável e decisões")
    tb(s, Inches(0.55), Inches(1.68), Inches(12.2), Inches(0.55),
       "Para cada variável: comparamos o valor máximo com o percentil 99 (razão máx/p99) e medimos a fração acima do limite IQR (Q3 + 1,5×IQR). Razão alta = pontos isolados suspeitos; fração alta = cauda longa estrutural (não é erro).",
       size=12.5, color=GREY)
    destaques = [
        ("Y (custo interno/km)", "máx 29,3 vs p99 1,50 (20×)", "cap no p99,5 na modelagem; custos negativos (estornos) excluídos"),
        ("n_os_acum", "máx 15.006 vs p99 77 (195×)", "anomalia de cadastro em pouquíssimas carretas — winsorizar/monitorar; árvores e imputação por mediana amortecem"),
        ("custo_acum_manutencao", "máx CAD 1,93 mi vs p99 32 mil (60×)", "mesma origem da anomalia acima — winsorizar/monitorar"),
        ("custo_medio_movel_3m", "máx 39 mil vs p99 1,5 mil (26×)", "trimestres excepcionais reais + anomalias — monitorar"),
        ("km_rodado_mes / km_acumulado", "máx 108 mil km/mês (8×) · 7,2 mi km (12×)", "outliers moderados após tratamento de resets — mantidos (robustez de árvore)"),
        ("ano_modelo, eixos, comprimento, franquia", "sem outliers relevantes", "nenhuma ação"),
    ]
    yy = 2.35
    tb(s, Inches(0.55), Inches(yy), Inches(3.3), Inches(0.3), "VARIÁVEL", size=12, color=ORANGE, bold=True)
    tb(s, Inches(3.95), Inches(yy), Inches(3.6), Inches(0.3), "EVIDÊNCIA", size=12, color=ORANGE, bold=True)
    tb(s, Inches(7.65), Inches(yy), Inches(5.2), Inches(0.3), "DECISÃO", size=12, color=ORANGE, bold=True)
    yy += 0.4
    for nome, evid, dec in destaques:
        tb(s, Inches(0.55), Inches(yy), Inches(3.3), Inches(0.65), nome, size=12.5, color=NAVY, bold=True)
        tb(s, Inches(3.95), Inches(yy), Inches(3.6), Inches(0.65), evid, size=11.5, color=DARK)
        tb(s, Inches(7.65), Inches(yy), Inches(5.2), Inches(0.7), dec, size=11.5, color=DARK)
        yy += 0.68
    tb(s, Inches(0.55), Inches(6.75), Inches(12.2), Inches(0.6),
       "Regra geral: caudas longas de custo são a natureza do fenômeno (não erro de dado) — mantidas, com cap apenas no alvo. Tabela completa: reports/tables/03d_diagnostico_outliers.csv.",
       size=12, color=NAVY, bold=True)

# sintese EDA
s = add_slide()
header(s, "EDA · síntese", "O que a EDA diz para a modelagem")
bullets(s, Inches(0.55), Inches(1.8), Inches(12.2), Inches(5.2), [
    (f"Zero-inflação define o problema: {stats['share_y_zero']*100:.0f}% dos meses sem custo — prever custo/km combina ocorrência + magnitude", 0, True),
    ("Cauda pesada nos positivos (assimetria alta): usar log1p / perda robusta; outliers extremos capados no p99,5", 0, False),
    ("Sinal fraco e distribuído: nenhuma X isolada passa de ρ ≈ 0,22 — ganho vem de interações e ensembles", 0, False),
    ("O histórico da carreta (OS acumuladas, custo acumulado, intervalo entre OS) é o bloco preditivo mais forte — H4 suportada", 0, True),
    ("Colinearidade a tratar: pares acumulados totais × preventivos; idade × ano_modelo", 0, False),
    ("Deflacionar é obrigatório: custos em CAD corrigidos pelo CPI Canadá (≈ +20% acumulado na janela); mesmo em valores reais o custo por km cresce ≈ 71% (2020→2025)", 0, True),
    ("km_rodado_mes é denominador do Y e feature: interpretação cuidadosa (relação mecânica)", 0, False),
])

# ============================================================ TECNICAS + MODELAGEM
divider("Bloco 2 · Pergunta 9", "Técnicas estatísticas e de Machine Learning",
        "Aplicadas sobre a população MAINT · split temporal treino/teste")

s = add_slide()
header(s, "Pergunta 9", "Técnicas utilizadas")
tb(s, Inches(0.55), Inches(1.75), Inches(6.0), Inches(0.4), "ESTATÍSTICA", size=14, color=ORANGE, bold=True)
bullets(s, Inches(0.55), Inches(2.15), Inches(6.0), Inches(2.6), [
    ("Correlação de Pearson e Spearman", 0, False),
    ("Análise de associação para categóricas (eta)", 0, False),
    ("Regressão linear simples e múltipla", 0, False),
    ("Ridge e Ridge com log1p no alvo", 0, False),
    ("Regressão polinomial (grau 2)", 0, False),
    ("Diagnóstico de multicolinearidade (VIF)", 0, False),
], size=14, gap=4)
tb(s, Inches(6.9), Inches(1.75), Inches(6.0), Inches(0.4), "MACHINE LEARNING", size=14, color=ORANGE, bold=True)
bullets(s, Inches(6.9), Inches(2.15), Inches(6.0), Inches(2.6), [
    ("Árvore de decisão para regressão", 0, False),
    ("Random Forest", 0, False),
    ("Gradient Boosting", 0, False),
    ("KNN (benchmark amostral)", 0, False),
    ("Modelo hurdle em duas partes (classificador de ocorrência × regressor de magnitude)", 0, False),
], size=14, gap=4)
tb(s, Inches(0.55), Inches(4.75), Inches(12.2), Inches(0.4), "AVALIAÇÃO E PREPARO", size=14, color=ORANGE, bold=True)
bullets(s, Inches(0.55), Inches(5.15), Inches(12.2), Inches(2.1), [
    ("Split temporal em palavras simples: o modelo estuda com o passado (2020–2024) e faz a prova com 2025, um período que ele nunca viu — exatamente como será usado na vida real", 0, True),
    ("Validação temporal expansiva (3 rodadas) · imputação mediana/moda · padronização · one-hot com agrupamento de categorias raras", 0, False),
    ("Métricas: R² (quanto da variação o modelo explica), RMSE e MAE (tamanho médio do erro, em CAD/km)", 0, False),
], size=14, gap=4)

# resultados modelos
if metricas is not None:
    s = add_slide()
    header(s, "Modelagem", "Comparação de modelos no teste temporal")
    m = metricas.copy()
    cols_show = [c for c in ["modelo", "r2", "rmse", "mae"] if c in m.columns]
    m = m[cols_show].sort_values("r2", ascending=False)
    yy = 1.9
    heads = ["MODELO", "R²", "RMSE", "MAE"][: len(cols_show)]
    xs = [Inches(0.7), Inches(7.2), Inches(9.0), Inches(10.9)]
    for j, htxt in enumerate(heads):
        tb(s, xs[j], Inches(yy), Inches(1.9), Inches(0.35), htxt, size=13, color=ORANGE, bold=True)
    yy += 0.42
    best = modelo_rec.iloc[0]["modelo_recomendado"] if modelo_rec is not None and "modelo_recomendado" in modelo_rec.columns else None
    for _, r in m.iterrows():
        is_best = best is not None and r[cols_show[0]] == best
        c = NAVY if is_best else DARK
        tb(s, xs[0], Inches(yy), Inches(6.3), Inches(0.38), str(r[cols_show[0]]) + ("  ◀ recomendado" if is_best else ""),
           size=13, color=c, bold=is_best)
        for j, colname in enumerate(cols_show[1:], 1):
            try:
                tb(s, xs[j], Inches(yy), Inches(1.8), Inches(0.38), fmt(float(r[colname]), 4), size=13, color=c, bold=is_best)
            except (ValueError, TypeError):
                tb(s, xs[j], Inches(yy), Inches(1.8), Inches(0.38), str(r[colname]), size=13, color=c)
        yy += 0.42
    tb(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(1.0),
       "Em palavras simples: R² = 0,086 significa que o modelo explica ~9% da variação do custo mês a mês. Parece pouco — e é mesmo modesto — mas é o suficiente para ORDENAR as carretas (quais tendem a custar mais) e apoiar orçamento; não serve para cravar o valor exato de um mês específico. O custo mensal tem muito acaso (um pneu fura ou não fura). Esses números são honestos: versões anteriores pareciam melhores porque \"colavam\" usando informação do futuro (vazamento), que foi removida.",
       size=12, color=NAVY, bold=True)

# importancia variaveis
if perm is not None and len(perm) > 0:
    s = add_slide()
    header(s, "Modelagem", "Fatores mais importantes (permutation importance, teste temporal)")
    pcols = perm.columns.tolist()
    name_c = pcols[0]
    val_c = [c for c in pcols if "import" in c.lower() or "mean" in c.lower()]
    val_c = val_c[0] if val_c else pcols[1]
    top = perm.nlargest(10, val_c)
    yy = 2.0
    maxv = float(top[val_c].max())
    for _, r in top.iterrows():
        w = max(0.15, 6.5 * float(r[val_c]) / maxv)
        bar = s.shapes.add_shape(1, Inches(4.3), Inches(yy), Inches(w), Inches(0.32))
        bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background(); bar.shadow.inherit = False
        tb(s, Inches(0.55), Inches(yy - 0.04), Inches(3.7), Inches(0.4), str(r[name_c])[:34], size=12, color=DARK)
        tb(s, Inches(4.35 + w), Inches(yy - 0.04), Inches(1.6), Inches(0.4), f"{float(r[val_c]):.4f}", size=11, color=GREY)
        yy += 0.44
    tb(s, Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.9),
       "Como ler (em palavras simples): para medir a importância de cada variável, nós a \"embaralhamos\" e vemos o quanto a previsão piora — quanto maior a barra, mais o modelo depende dela. Confirma a EDA: o histórico operacional e a intensidade de uso dominam; atributos fixos (ano, eixos) pesam menos.",
       size=12, color=NAVY, bold=True)

# hipoteses x evidencias
s = add_slide()
header(s, "Síntese", "Hipóteses × evidências")
hips = [
    ("H1 · duração de contrato ⇒ custo", "NÃO SUPORTADA", "correlação direta ≈ 0 (ρ = +0,02) no grão mensal"),
    ("H2 · idade ⇒ custo", "PARCIAL", "efeito direto fraco (ρ = +0,04); idade opera via histórico acumulado (colinearidade)"),
    ("H3 · quilometragem ⇒ custo", "PARCIAL", "km_acumulado ρ = +0,16 (estrutural); km mensal tem relação mecânica com o denominador"),
    ("H4 · histórico prevê custo futuro", "SUPORTADA", "n_os_acum ρ = +0,22, custo_acum ρ = +0,20, intervalo_medio_os ρ = −0,19 — bloco mais forte do ranking"),
    ("H5 · operação e contrato influenciam", "PARCIAL", "região é a categórica mais forte (eta = 0,084); reefer e montadora deslocam medianas; efeito contratual fraco"),
]
yy = 1.9
for nome, status, evid in hips:
    color = RGBColor(0x2E, 0x7D, 0x32) if status == "SUPORTADA" else (RGBColor(0xC6, 0x28, 0x28) if "NÃO" in status else ORANGE)
    tb(s, Inches(0.55), Inches(yy), Inches(4.6), Inches(0.7), nome, size=14, color=NAVY, bold=True)
    tb(s, Inches(5.3), Inches(yy), Inches(1.9), Inches(0.7), status, size=13, color=color, bold=True)
    tb(s, Inches(7.3), Inches(yy), Inches(5.5), Inches(0.9), evid, size=12, color=DARK)
    yy += 0.95

# ============================================================ BLOCO 3
divider("Bloco 3", "Referencial teórico e metodologia",
        "Perguntas 10–11 da agenda da disciplina")

s = add_slide()
header(s, "Pergunta 10", "Referencial teórico — síntese")
bullets(s, Inches(0.55), Inches(1.8), Inches(12.2), Inches(5.0), [
    ("Katreddi et al. (2023a) — ensembles (Super Learner) sobre variáveis operacionais alcançam alto desempenho em custo/milha ⇒ prioridade a modelos de árvore/ensemble no projeto", 0, False),
    ("Katreddi et al. (2023b) — Mixed Effects Random Forest captura heterogeneidade entre grupos de veículos ⇒ variáveis de agrupamento (montadora, classe, contrato) entram como features; MERF é extensão futura", 0, False),
    ("Sun et al. (2024) — histórico de manutenção sustenta previsão de custos (via confiabilidade/Weibull) ⇒ valida o uso de features de histórico acumulado (nosso bloco preditivo mais forte)", 0, False),
    ("Adekitan et al. (2018) — ANN com R = 0,766 mostra sinal em variáveis operacionais, com desempenho modesto ⇒ referência histórica; redes neurais não priorizadas", 0, False),
    ("Convergência: variáveis operacionais e de histórico > atributos estáticos; métodos de árvore > lineares; heterogeneidade entre grupos importa", 0, True),
])

s = add_slide()
header(s, "Pergunta 11", "Metodologia — pipeline em 7 notebooks reprodutíveis")
metodo = [
    ("00 · Contexto e inventário", "documentação do problema, inventário das 7 bases, validação do modelo estrela, riscos iniciais"),
    ("01 · Qualidade e integridade", "chaves, duplicidades, ausentes, custos negativos, odômetro (resets/saltos), regras de tratamento"),
    ("02 · Base analítica mensal", "agregação carreta × mês, prorrateio de km, piso 500 km/mês, features históricas defasadas (anti-vazamento)"),
    ("03 · EDA e hipóteses", "protocolo variável a variável, correlações, associação de categóricas, VIF, teste das 5 hipóteses"),
    ("04 · Deflação (CPI Canadá)", "custos CAD convertidos a valores de dez/2025 via CPI all-items StatCan (correção do IPCA)"),
    ("05 · Modelagem preditiva", "split temporal, 9 modelos + hurdle, métricas R²/RMSE/MAE, permutation importance, erro por segmento"),
    ("06 · Resultados", "síntese executiva, hipóteses finais, recomendações de negócio e limitações"),
]
yy = 1.85
for nome, desc_ in metodo:
    tb(s, Inches(0.55), Inches(yy), Inches(3.7), Inches(0.65), nome, size=14, color=NAVY, bold=True)
    tb(s, Inches(4.4), Inches(yy), Inches(8.3), Inches(0.7), desc_, size=12.5, color=DARK)
    yy += 0.72

# limitacoes
s = add_slide()
header(s, "Transparência metodológica", "Limitações e cuidados")
bullets(s, Inches(0.55), Inches(1.8), Inches(12.2), Inches(5.2), [
    ("Zero-inflação: 67% dos meses sem custo — desempenho preditivo pontual é modesto; o modelo apoia planejamento, não substitui validação operacional", 0, False),
    ("km_rodado_mes é denominador do alvo e feature — interpretação cuidadosa", 0, False),
    ("Resets e saltos de odômetro tratados, mas o prorrateio de km é aproximação", 0, False),
    ("Custos negativos (estornos/ajustes contábeis) excluídos da modelagem", 0, False),
    ("GPS com cobertura parcial (set–dez/2025) — região derivada da província/local da OS", 0, False),
    ("Contratos vigentes têm duração censurada (referência 2025-12)", 0, False),
    ("Vazamento temporal removido (meses_desde_ultima_os e km_acumulado defasados): métricas atuais são menores porém defensáveis", 0, True),
    ("Extensão futura: Mixed-Effects Random Forest / modelos hierárquicos e modelos zero-inflados", 0, False),
])

# encerramento
s = add_slide(NAVY)
tb(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.0), "Obrigado!", size=44, color=WHITE, bold=True)
tb(s, Inches(0.9), Inches(3.8), Inches(11.5), Inches(1.4),
   "Previsão de Custos de Manutenção de Carretas\nQuatro Norte Consulting · MBA FGV · 2026",
   size=18, color=LIGHT)
bar = s.shapes.add_shape(1, Inches(0.9), Inches(3.55), Inches(2.6), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background(); bar.shadow.inherit = False

prs.save(OUT)
print(f"OK: {OUT} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
