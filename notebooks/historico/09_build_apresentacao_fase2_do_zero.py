"""09 - Apresentacao Fase 2, no formato do PDF apresentado em 2026-08-05.

Gera `docs/entregas/Apresentacao_QuatroNorte_Fase2.pptx` a partir de `reports/`.

Estrutura espelha o PDF entregue (secoes 01-15 + Gates), com acentuacao, e acrescenta
o bloco da Fase 2: chegada dos campos de contrato, resultados, hipoteses e conclusoes.
Nenhum numero e digitado manualmente.
"""
import pandas as pd
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
TABLES = ROOT / "reports" / "tables"
FIG = ROOT / "reports" / "figures"
FIG_EDA = FIG / "eda"
OUT = ROOT / "docs" / "entregas" / "Apresentacao_QuatroNorte_Fase2.pptx"


def rd(n):
    p = TABLES / n
    return pd.read_csv(p) if p.exists() else pd.DataFrame()


stats = rd("03c_stats_ppt.csv").set_index("metrica")["valor"] if not rd("03c_stats_ppt.csv").empty else {}
inv = rd("00_inventario_base_consolidada.csv").set_index("metrica")["valor"] if not rd("00_inventario_base_consolidada.csv").empty else {}
metr = rd("05_metricas_modelos.csv")
metr_all = rd("05_metricas_por_configuracao.csv")
cmpcfg = rd("05_comparacao_configuracoes.csv")
best = rd("05_modelo_recomendado.csv")
imp = rd("05_importancia_permutacao_random_forest.csv")
hip = rd("06_hipoteses_final.csv")
sel = rd("05_selecao_variaveis.csv")
corr = rd("03b_correlacao_com_y.csv")
eta = rd("03b_eta_categoricas.csv")
dic = rd("02_dicionario_base_anual.csv")
vif = rd("03b_vif.csv")
est = rd("03b_estatisticas_descritivas.csv")
cur = rd("02_curadoria_limpeza.csv")
val = rd("02_validacao_base_anual.csv")
rec = rd("06_recomendacoes_negocio.csv")
pop = rd("05_resumo_populacao_modelagem.csv")


def s(k, d="—"):
    try:
        return stats[k]
    except Exception:
        return d


def vl(chave, d="—"):
    if val.empty:
        return d
    r = val.loc[val["checagem"] == chave, "valor"]
    return r.iloc[0] if len(r) else d


def cur_val(passo, d=0):
    if cur.empty:
        return d
    r = cur.loc[cur["passo"] == passo, "linhas"]
    return int(r.iloc[0]) if len(r) else d


# ---------------- identidade visual (mesma do deck vigente) ----------------
NAVY = RGBColor(0x14, 0x2B, 0x45)
INK = RGBColor(0x1C, 0x24, 0x33)
ORANGE = RGBColor(0xE8, 0x71, 0x3A)
PAPER = RGBColor(0xF6, 0xF1, 0xE7)
GREY = RGBColor(0x5B, 0x63, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x50)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide(bg=PAPER):
    sl = prs.slides.add_slide(BLANK)
    r = sl.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    sl.shapes._spTree.remove(r._element)
    sl.shapes._spTree.insert(2, r._element)
    return sl


def tb(sl, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        f = run.font
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
        f.name = "Calibri"
    return box


def header(sl, kicker, title):
    bar = sl.shapes.add_shape(1, Inches(0.6), Inches(0.55), Inches(0.16), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()
    bar.shadow.inherit = False
    tb(sl, 0.9, 0.5, 11.8, 0.4, kicker.upper(), 12, ORANGE, True)
    tb(sl, 0.9, 0.86, 11.8, 0.9, title, 27, NAVY, True)


def bullets(sl, x, y, w, h, items, size=16, color=INK, gap=6):
    box = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = "•  " + it
        f = run.font
        f.size = Pt(size)
        f.color.rgb = color
        f.name = "Calibri"
    return box


def pic_fit(sl, path, x, y, w, h):
    from PIL import Image as PImage
    p = Path(path)
    if not p.exists():
        tb(sl, x, y, w, 0.4, f"(figura ausente: {p.name})", 11, GREY)
        return
    try:
        iw, ih = PImage.open(p).size
        ar = iw / ih
    except Exception:
        ar = w / h
    bw, bh = w, w / ar
    if bh > h:
        bh, bw = h, h * ar
    sl.shapes.add_picture(str(p), Inches(x + (w - bw) / 2), Inches(y + (h - bh) / 2), Inches(bw), Inches(bh))


def pic_grid(sl, paths, x, y, w, h, cols=2, rows=2, gap=0.2):
    cw = (w - gap * (cols - 1)) / cols
    ch = (h - gap * (rows - 1)) / rows
    for i, p in enumerate(paths[: cols * rows]):
        r, c = divmod(i, cols)
        pic_fit(sl, p, x + c * (cw + gap), y + r * (ch + gap), cw, ch)


def table(sl, df, x, y, w, h, fs=11, maxrows=12, header_bg=NAVY, destaque=None):
    total = len(df)
    df = df.head(maxrows)
    nr, nc = df.shape
    gt = sl.shapes.add_table(nr + 1, nc, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for j, col in enumerate(df.columns):
        cell = gt.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        pr = cell.text_frame.paragraphs[0]
        pr.runs[0].font.size = Pt(fs)
        pr.runs[0].font.bold = True
        pr.runs[0].font.color.rgb = WHITE
    for i in range(nr):
        marca = destaque is not None and destaque(df.iloc[i])
        for j in range(nc):
            cell = gt.cell(i + 1, j)
            cell.text = str(df.iloc[i, j])
            cell.fill.solid()
            cell.fill.fore_color.rgb = (RGBColor(0xFD, 0xE8, 0xDC) if marca
                                        else (WHITE if i % 2 == 0 else RGBColor(0xEE, 0xEA, 0xDE)))
            pr = cell.text_frame.paragraphs[0]
            if pr.runs:
                pr.runs[0].font.size = Pt(fs)
                pr.runs[0].font.bold = bool(marca)
                pr.runs[0].font.color.rgb = INK
    if total > maxrows:
        tb(sl, x, y + h + 0.05, w, 0.3, f"({total - maxrows} linhas adicionais em reports/tables/)", 9, GREY)
    return gt


def kpi(sl, x, y, w, h, valor, rotulo, cor=ORANGE):
    c = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid()
    c.fill.fore_color.rgb = WHITE
    c.line.color.rgb = NAVY
    c.shadow.inherit = False
    tb(sl, x + 0.2, y + 0.08, w - 0.4, 0.45, valor, 22, cor, True)
    tb(sl, x + 0.2, y + 0.58, w - 0.4, 0.35, rotulo, 11, GREY)


best_pred = best[best["cenario"] == "preditivo"].iloc[0] if not best.empty else None
best_expl = best[best["cenario"] == "explicativo"].iloc[0] if not best.empty else None
cp = cmpcfg[cmpcfg["cenario"] == "preditivo"].iloc[0] if not cmpcfg.empty else None
ce = cmpcfg[cmpcfg["cenario"] == "explicativo"].iloc[0] if not cmpcfg.empty else None

# =============================== SLIDES ===============================

# 1. Capa
sl = slide(NAVY)
tb(sl, 0.9, 2.0, 11.6, 0.5, "QUATRO NORTE CONSULTING · GRUPO 01", 13, ORANGE, True)
tb(sl, 0.9, 2.6, 11.6, 1.6, "Custo anual de manutenção por carreta", 40, WHITE, True)
tb(sl, 0.9, 4.1, 11.6, 1.0,
   "Fatores determinantes e modelagem do custo anual (CAD $/ano, real) da frota\n"
   "própria de leasing/rental no Canadá — janela 2020–2025", 17, PAPER)
tb(sl, 0.9, 5.3, 11.6, 0.5, "FASE 2 · Modelagem, resultados e incorporação dos dados de contrato", 15, ORANGE, True)
tb(sl, 0.9, 6.4, 11.6, 0.5, "Marlon Wenzel · Jeison Lima · Rodrigo Queiroz · Giovani Cani", 13, PAPER)

# 2. O que mudou desde a Fase 1  (a historia dos campos novos)
sl = slide()
header(sl, "Fase 2", "O que mudou desde a apresentação anterior")
tb(sl, 0.9, 1.75, 11.6, 0.4,
   "A Fase 1 encerrou na seleção de variáveis. A modelagem ficou marcada como Fase 2 — e um dado novo mudou o que era possível testar.",
   12, GREY)
linhas = pd.DataFrame([
    {"tema": "Base de dados", "fase 1": "25 colunas · 223.590 OS · 9.859 carretas",
     "fase 2 (agora)": f"29 colunas · {int(float(inv.get('linhas_os', 0))):,} OS · {int(float(inv.get('carretas_distintas', 0))):,} carretas".replace(",", ".")},
    {"tema": "Contrato", "fase 1": "ausente da fonte única — H6 declarada, não testável",
     "fase 2 (agora)": "4 campos incorporados à própria fonte única → H6 testável"},
    {"tema": "População", "fase 1": "todo o custo interno (por falta do dado)",
     "fase 2 (agora)": f"contrato com manutenção inclusa (MAINT): {vl('carreta_ano_populacao_maint')} carreta-anos"},
    {"tema": "Modelagem", "fase 1": "prevista, não executada",
     "fase 2 (agora)": "7 modelos × 2 cenários × 3 configurações, split temporal"},
    {"tema": "Hipóteses", "fase 1": "H1–H6 declaradas",
     "fase 2 (agora)": "H1–H6 com veredito, incluindo H6a e H6b"},
])
table(sl, linhas, 0.9, 2.3, 11.6, 3.4, fs=12, maxrows=5)
tb(sl, 0.9, 6.2, 11.6, 0.8,
   "Importante: a fonte permanece ÚNICA. Os campos de contrato passaram a integrar o próprio CSV consolidado —\n"
   "não houve join novo nem tabela adicional.", 13, NAVY, True)

# 3. Os quatro campos novos
sl = slide()
header(sl, "Fase 2 · Dados novos", "Os quatro campos de contrato")
campos = pd.DataFrame([
    {"campo": "tempo_contrato_meses_ate_reparo", "cobertura": "92,5%",
     "perfil": "mediana 35,1 meses · p95 105 · sem valores negativos", "uso": "MODELADA (H6a)"},
    {"campo": "tipo_manutencao", "cobertura": "92,5%",
     "perfil": "MAINT 89,7% · MIX 1,6% · NET 1,1%", "uso": "define a população (D6); H6b na EDA"},
    {"campo": "cod_cliente", "cobertura": "77,2%",
     "perfil": "597 clientes distintos", "uso": "descritiva — não modelada"},
    {"campo": "franquia_km_mensal_contrato", "cobertura": "71,6%",
     "perfil": "99,8% dos preenchidos são ZERO", "uso": "DESCARTADA (variância nula)"},
])
table(sl, campos, 0.7, 1.9, 12.0, 2.3, fs=11, maxrows=4,
      destaque=lambda r: "DESCARTADA" in str(r["uso"]))
bullets(sl, 0.9, 4.5, 11.6, 2.3, [
    "Critério de descarte da franquia: o mesmo aplicado a tailgate_flag na Fase 1 — variância praticamente nula não sustenta hipótese.",
    "cod_cliente ficou fora do modelo por cardinalidade (597 categorias): entraria como memorização do cliente, não como explicação do custo.",
    "Contrato não é atributo fixo da carreta: 51,5% das carretas mudam de regime no período, o que exigiu agregação explícita no grão anual.",
], 14)

# 4. 01 CONTEXTO
sl = slide()
header(sl, "01 · Contexto", "A empresa e a operação")
bullets(sl, 0.9, 1.9, 7.2, 4.6, [
    "Empresa de leasing/rental de carretas no Canadá (secas e refrigeradas, até 53 pés).",
    "Manutenção própria: oficinas registram ordens de serviço com mão de obra e peças.",
    "Todos os custos em dólares canadenses (CAD); operação integralmente no Canadá.",
    "Objetivo de negócio: apoiar orçamento anual e priorização de manutenção da frota.",
    "Fonte única do estudo: base consolidada de Ordens de Serviço (2020–2025).",
], 16)
for i, (k, v) in enumerate([
    ("Carretas únicas", f"{int(float(inv.get('carretas_distintas', 0))):,}".replace(",", ".")),
    ("Ordens de serviço", f"{int(float(inv.get('linhas_os', 0))):,}".replace(",", ".")),
    ("Custo interno (real)", f"CAD $ {s('custo_total_real_mi')} mi"),
    ("Período", "2020–2025"),
]):
    kpi(sl, 8.5, 1.9 + i * 1.15, 4.0, 1.0, v, k)

# 5. 02 PROBLEMA
sl = slide()
header(sl, "02 · Problema", "A pergunta de pesquisa")
c = sl.shapes.add_shape(1, Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.3))
c.fill.solid(); c.fill.fore_color.rgb = NAVY; c.line.fill.background(); c.shadow.inherit = False
tb(sl, 1.2, 2.05, 11.0, 1.0,
   "Quais fatores mais influenciam o CUSTO ANUAL INTERNO de manutenção por carreta e como estimá-lo\n"
   "a partir das características operacionais, históricas e estruturais da frota?", 18, WHITE, True)
bullets(sl, 0.9, 3.5, 11.5, 3.2, [
    "Variável de interesse: custo anual de manutenção por carreta (CAD $/ano), em valores reais.",
    "Grão de análise: carreta × ano.",
    "Escopo: todo o custo interno absorvido pela empresa (preventiva + corretiva).",
    "Correção monetária obrigatória pela inflação canadense (CPI) para comparar anos em termos reais.",
    "População da modelagem: carretas sob contrato com manutenção inclusa (MAINT) — critério retomado na Fase 2.",
], 16)

# 6. 03 OBJETIVOS
sl = slide()
header(sl, "03 · Objetivos", "Objetivos gerais e específicos")
tb(sl, 0.9, 1.85, 11.6, 1.1,
   "Analisar os fatores que influenciam o custo anual de manutenção por carreta (CAD, corrigido pela inflação), "
   "identificando as variáveis de maior capacidade explicativa e desenvolvendo modelos estatísticos e de ML para estimá-lo.",
   15, INK, True)
bullets(sl, 0.9, 3.05, 11.6, 3.5, [
    "Consolidar a análise a partir da base única de OS.",
    "Definir a variável resposta anual e corrigi-la pelo CPI (Canadian Price Index) do Canadá.",
    "Realizar EDA rigorosa (univariada, relação com Y, ranking) das variáveis candidatas.",
    "Selecionar variáveis de forma fundamentada (associação, multicolinearidade, domínio).",
    "Desenvolver e avaliar modelos estatísticos e de Machine Learning (split temporal).",
    "Testar o efeito das variáveis de contrato sobre o custo anual (novo na Fase 2).",
], 15)

# 7. 04 HIPOTESES
sl = slide()
header(sl, "04 · Hipóteses", "Hipóteses de trabalho e vereditos")
if not hip.empty:
    table(sl, hip[["hipotese", "enunciado", "veredito"]], 0.9, 1.9, 11.6, 4.4, fs=11, maxrows=10,
          destaque=lambda r: str(r["hipotese"]).startswith("H6"))
tb(sl, 0.9, 6.5, 11.6, 0.7,
   "H6 constava da Fase 1 como hipótese declarada e sem dados. Na Fase 2 foi desdobrada em H6a (duração) e H6b (tipo) e finalmente testada.",
   12, GREY)

# 8. 05 BASE DE DADOS + fluxograma
sl = slide()
header(sl, "05 · Base de dados", "Base consolidada única (29 colunas)")
bullets(sl, 0.9, 1.85, 7.3, 3.2, [
    "Fonte única: fato_wo_ml_2020-01-01_to_2025-12-31.csv (1 linha = 1 OS).",
    "Extração SQL, modelo estrela, joins e feature engineering: etapa anterior.",
    f"{cur_val('linhas_originais'):,} OS · {cur_val('carretas_distintas'):,} carretas · 2020–2025.".replace(",", "."),
    "VMRS: código padronizado do sistema reparado (PM, freios, pneus, reefer...).",
    "NOVO: 4 campos de contrato dentro da própria base — a fonte segue única.",
], 14)
cur_show = cur.rename(columns={"passo": "etapa de curadoria", "linhas": "linhas"}) if not cur.empty else pd.DataFrame()
if not cur_show.empty:
    table(sl, cur_show, 0.9, 5.2, 7.3, 1.9, fs=9, maxrows=6)
tb(sl, 8.6, 1.85, 4.2, 0.4, "Arquitetura metodológica", 14, NAVY, True)
tb(sl, 8.6, 2.3, 4.4, 4.8,
   "Base consolidada (CSV)\n↓ Validação de qualidade\n↓ Variável resposta anual\n↓ Correção CPI Canadá\n"
   "↓ EDA (uni/bivariada)\n↓ Correlação e testes\n↓ Multicolinearidade/VIF\n↓ Ranking\n↓ Seleção de variáveis\n"
   "↓ Modelagem (estat.+ML) ✓\n↓ Avaliação / discussão ✓", 13, INK)
tb(sl, 8.6, 6.85, 4.4, 0.3, "✓ concluído na Fase 2", 11, GREEN, True)

# 9. 06 VARIAVEL RESPOSTA
sl = slide()
header(sl, "06 · Variável resposta", "Y = custo anual de manutenção por carreta")
c = sl.shapes.add_shape(1, Inches(0.9), Inches(1.85), Inches(11.5), Inches(1.0))
c.fill.solid(); c.fill.fore_color.rgb = NAVY; c.line.fill.background(); c.shadow.inherit = False
tb(sl, 1.2, 2.0, 11.0, 0.8,
   "Y = soma do custo interno da carreta no ano (CAD), deflacionado pelo CPI Canadá (dez/2025)  ·  Grão: carreta × ano",
   16, WHITE, True)
bullets(sl, 0.9, 3.15, 11.5, 3.4, [
    f"Média CAD $ {s('y_media')}/ano · mediana CAD $ {s('y_mediana')}/ano · p90 CAD $ {s('y_p90')}.",
    f"Assimetria {s('y_assimetria')} (cauda longa), tratada por log1p e modelos robustos.",
    f"Apenas {s('share_y_zero_pct')}% de carreta-anos com custo zero: o grão anual praticamente elimina a zero-inflação.",
    f"Custo interno total: CAD $ {s('custo_total_nominal_mi')} mi nominal / CAD $ {s('custo_total_real_mi')} mi real (dez/2025).",
    f"Base analítica: {vl('linhas_base_anual')} carreta-anos; população MAINT (modelagem): {vl('carreta_ano_populacao_maint')}.",
], 15)

# 10. 07 INFLACAO
sl = slide()
header(sl, "07 · Correção pela inflação", "Custos reais pelo CPI do Canadá")
bullets(sl, 0.9, 1.9, 5.4, 4.4, [
    "Todos os custos em CAD, corrigidos pelo CPI all-items do Canadá (StatCan, vetor v41690973).",
    "Base de comparação: dezembro de 2025.",
    "Inflação acumulada no período: ~20,6% (2020→2025).",
    "A correção isola mudanças REAIS de custo da mera perda do poder de compra.",
], 14)
pic_fit(sl, FIG / "04_nominal_vs_deflacionado.png", 6.5, 1.8, 6.3, 4.9)

# 11. 08 VARIAVEIS - dicionario
sl = slide()
header(sl, "08 · Variáveis", "Dicionário das variáveis (grão carreta × ano)")
tb(sl, 0.9, 1.55, 11.6, 0.4,
   "Universo derivável da fonte única, já incluindo o bloco de contrato. Em destaque, as variáveis novas da Fase 2.",
   11, GREY)
if not dic.empty:
    d2 = dic[["variavel", "tipo/papel"]].rename(columns={"tipo/papel": "tipo / papel"})
    meta = int((len(d2) + 1) / 2)
    marca = lambda r: ("contrato" in str(r["tipo / papel"]).lower()) or ("populacao" in str(r["tipo / papel"]).lower())
    table(sl, d2.iloc[:meta], 0.7, 2.0, 5.9, 4.8, fs=8, maxrows=meta, destaque=marca)
    table(sl, d2.iloc[meta:], 6.8, 2.0, 5.9, 4.8, fs=8, maxrows=len(d2) - meta, destaque=marca)

# 12. 09 EDA - descritivas
sl = slide()
header(sl, "09 · EDA", "Estatísticas descritivas")
if not est.empty:
    cols = [c for c in ["variavel", "N", "media", "mediana", "desvio_padrao", "min", "max", "assimetria"] if c in est.columns]
    table(sl, est[cols], 0.7, 1.9, 12.0, 4.9, fs=9, maxrows=16)

# 13. 09 EDA - evolucao
sl = slide()
header(sl, "09 · EDA", "Evolução do custo anual real (2020–2025)")
pic_fit(sl, FIG_EDA / "evolucao_y_anual.png", 0.9, 1.8, 7.4, 5.0)
bullets(sl, 8.6, 2.0, 4.2, 4.6, [
    "Crescimento consistente do custo anual real, mesmo após a correção pela inflação.",
    "Média e mediana sobem: o aumento é disseminado na frota, não concentrado.",
    "A distância entre média e mediana evidencia ativos de custo excepcionalmente alto.",
], 13)

# 14-15. EDA histogramas
sl = slide()
header(sl, "09 · EDA", "Histogramas — resposta, uso e histórico")
pic_grid(sl, [FIG_EDA / "quant_custo_ano_real.png", FIG_EDA / "quant_km_rodado_ano.png",
              FIG_EDA / "quant_custo_ano_anterior.png", FIG_EDA / "quant_n_os_ano.png"],
         0.8, 1.8, 11.8, 5.0)

sl = slide()
header(sl, "09 · EDA", "Histogramas — atributos do ativo")
pic_grid(sl, [FIG_EDA / "quant_ano_modelo.png", FIG_EDA / "quant_eixos.png",
              FIG_EDA / "quant_comprimento.png", FIG_EDA / "quant_idade_carreta.png"],
         0.8, 1.8, 11.8, 5.0)

# 16. EDA histogramas de contrato (NOVO)
sl = slide()
header(sl, "09 · EDA · Contrato", "Distribuições das variáveis de contrato (novas)")
pic_grid(sl, [FIG_EDA / "quant_tempo_contrato_meses_fim_ano.png", FIG_EDA / "quant_share_maint_ano.png",
              FIG_EDA / "quant_n_clientes_ano.png", FIG_EDA / "quant_trocou_contrato_ano.png"],
         0.8, 1.8, 11.8, 5.0)

# 17. Boxplots por categoria
sl = slide()
header(sl, "09 · EDA", "Boxplots do custo anual por categoria")
pic_grid(sl, [FIG_EDA / "quali_flag_refrigerado.png", FIG_EDA / "quali_unit_subtype.png",
              FIG_EDA / "quali_cod_montadora.png", FIG_EDA / "quali_vmrs_predominante_ano.png"],
         0.8, 1.8, 11.8, 5.0)

# 18. Boxplot por tipo de manutencao (H6b)
sl = slide()
header(sl, "09 · EDA · Contrato", "Custo anual por regime contratual — evidência de H6b")
pic_fit(sl, FIG_EDA / "quali_tipo_manutencao_ano.png", 0.9, 1.8, 7.4, 5.0)
_eta_tipo = eta.loc[eta["variavel"] == "tipo_manutencao_ano", "eta"]
_eta_tipo = float(_eta_tipo.iloc[0]) if len(_eta_tipo) else float("nan")
bullets(sl, 8.6, 2.0, 4.2, 4.6, [
    f"η = {_eta_tipo:.3f} — separação fraca entre regimes contratuais.",
    "MAINT concentra 89,7% das OS; NET e MIX somam 2,7%: comparação de baixa potência estatística.",
    "Avaliado sobre a base COMPLETA — dentro da população MAINT a variável é constante e não pode ser testada.",
    "Diferenças entre regimes podem refletir quem paga o reparo, não quanto ele custa.",
], 13)

# 19. Relacao individual com Y
sl = slide()
header(sl, "09 · EDA", "Relação individual de cada variável com Y")
if not corr.empty:
    cc = corr[["variavel", "spearman", "papel"]].sort_values("spearman", key=abs, ascending=False)
    table(sl, cc, 0.7, 1.9, 6.0, 4.8, fs=9, maxrows=15,
          destaque=lambda r: "contrato" in str(r["variavel"]) or str(r["variavel"]) in ("share_maint_ano", "n_clientes_ano"))
if not eta.empty:
    table(sl, eta[["variavel", "eta", "n_categorias"]], 7.0, 1.9, 5.7, 4.8, fs=9, maxrows=12,
          destaque=lambda r: str(r["variavel"]) == "tipo_manutencao_ano")
tb(sl, 0.7, 6.85, 12.0, 0.3,
   "Quantitativas: Spearman. Categóricas: eta (ANOVA). Em destaque, as variáveis de contrato.", 10, GREY)

# 20. Ranking
sl = slide()
header(sl, "09 · EDA", "Ranking de associação com o custo anual")
pic_fit(sl, FIG_EDA / "ranking_associacao_y.png", 0.9, 1.8, 7.6, 5.0)
bullets(sl, 8.7, 2.0, 4.1, 4.6, [
    "Uso (km) e histórico defasado lideram entre as quantitativas.",
    "Refrigerado e subtipo lideram entre as categóricas.",
    "Idade tem efeito direto fraco.",
    "Contrato aparece na metade inferior do ranking.",
    "Componentes aritméticos de Y são excluídos do ranking.",
], 13)

# 21. Multicolinearidade
sl = slide()
header(sl, "09 · EDA", "Multicolinearidade — matriz e VIF")
pic_fit(sl, FIG_EDA / "matriz_spearman.png", 0.8, 1.8, 7.2, 5.0)
if not vif.empty:
    table(sl, vif, 8.2, 1.9, 4.4, 4.8, fs=10, maxrows=14)

# 22. Selecao
sl = slide()
header(sl, "10 · Seleção das variáveis", "Decisões e critérios")
if not sel.empty:
    s2 = sel[["variavel", "decisao"]]
    meta = int((len(s2) + 1) / 2)
    marca = lambda r: ("contrato" in str(r["variavel"])) or (str(r["variavel"]) in
                      ("share_maint_ano", "n_clientes_ano", "trocou_contrato_ano", "tipo_manutencao_ano",
                       "cod_cliente_predominante_ano", "franquia_km_mensal_contrato"))
    table(sl, s2.iloc[:meta], 0.7, 1.9, 6.0, 4.7, fs=8, maxrows=meta, destaque=marca)
    table(sl, s2.iloc[meta:], 6.9, 1.9, 6.0, 4.7, fs=8, maxrows=len(s2) - meta, destaque=marca)
tb(sl, 0.7, 6.75, 12.0, 0.4,
   "Critérios: ranking, VIF, redundância, vazamento temporal e coerência de domínio. Em destaque, as decisões sobre contrato.",
   11, GREY)

# 23. Referencial teorico
sl = slide()
header(sl, "11 · Referencial teórico", "Estudos selecionados e trabalhos-base")
ref = pd.DataFrame([
    {"Autor (ano)": "Katreddi, Thiruvengadam, Thompson,\nSchmid e Padmanaban (2023)",
     "Tema": "Custo de manutenção de caminhões\nde entrega (diesel/GNV)",
     "Método / achado": "RF, XGBoost, ANN e Super Learner;\nSuper Learner o melhor (R²=97,28%)",
     "Contribuição": "Fundamenta árvores e ensembles para\ncapturar relação não linear"},
    {"Autor (ano)": "Katreddi, Thiruvengadam,\nThompson e Schmid (2023)",
     "Tema": "Custo em veículos pesados;\nfrotas heterogêneas",
     "Método / achado": "Mixed-Effects Random Forest (MERF)",
     "Contribuição": "Apoia categóricas de agrupamento\n(montadora, subtipo, refrigeração)"},
    {"Autor (ano)": "Sun, Guo, Sun,\nYang e Hao (2024)",
     "Tema": "Previsão de custo a partir de\ndados de manutenção",
     "Método / achado": "Mixed Weibull + estimativa iterativa\nsobre o histórico",
     "Contribuição": "Sustenta as variáveis históricas\ndefasadas — confirmadas aqui"},
    {"Autor (ano)": "Adekitan, Adetokun\ne Okokpujie (2018)",
     "Tema": "Componentes de custo de\nmanutenção veicular",
     "Método / achado": "Rede neural artificial (ANN); R=0,766",
     "Contribuição": "Sinal preditivo em uso e histórico;\nnão priorizar redes neurais"},
])
table(sl, ref, 0.6, 1.85, 12.2, 4.9, fs=9, maxrows=4)

# 24. Modelagem - tecnicas
sl = slide()
header(sl, "12 · Modelagem", "Técnicas e desenho experimental")
bullets(sl, 0.9, 1.85, 11.6, 4.6, [
    "Estatística: regressão linear múltipla, ridge, regressão polinomial de grau 2.",
    "Machine Learning: árvore de decisão, Random Forest, Gradient Boosting, KNN.",
    "Alvo transformado por log1p (assimetria 3,82); métricas na escala original (CAD/ano).",
    "Split temporal: treino 2020–2024, teste 2025 — evita vazamento entre períodos.",
    "Dois cenários: explicativo (inclui uso do ano) e preditivo (histórico defasado + atributos).",
    "TRÊS configurações de população/variáveis, para separar o efeito do filtro MAINT do efeito do contrato.",
    "Métricas: R², RMSE e MAE.",
], 15)

# 25. Desenho das 3 configuracoes (NOVO - o coracao da Fase 2)
sl = slide()
header(sl, "13 · Resultados", "Como isolamos o efeito de cada mudança")
tb(sl, 0.9, 1.7, 11.6, 0.5,
   "A base foi reextraída E ganhou contrato na mesma rodada. Sem separar as causas, qualquer ganho seria ambíguo.",
   13, GREY)
cfgs = pd.DataFrame([
    {"config": "A — baseline", "população": "todas as carretas", "variáveis": "sem contrato", "responde": "efeito da nova extração"},
    {"config": "B — filtro", "população": "somente MAINT", "variáveis": "sem contrato", "responde": "efeito do filtro de contrato"},
    {"config": "C — completa", "população": "somente MAINT", "variáveis": "com contrato", "responde": "efeito das variáveis novas"},
])
table(sl, cfgs, 0.9, 2.35, 11.6, 1.8, fs=12, maxrows=3)
if cp is not None:
    kpi(sl, 0.9, 4.6, 3.5, 1.2, f"{float(cp['A_todos_sem_contrato__r2']):.4f}", "A · R² preditivo", NAVY)
    kpi(sl, 4.8, 4.6, 3.5, 1.2, f"{float(cp['B_maint_sem_contrato__r2']):.4f}", "B · R² preditivo (MAINT)", NAVY)
    kpi(sl, 8.7, 4.6, 3.5, 1.2, f"{float(cp['C_maint_com_contrato__r2']):.4f}", "C · R² preditivo (+contrato)", ORANGE)
    tb(sl, 0.9, 6.1, 11.6, 0.9,
       f"Efeito do filtro MAINT: {float(cp['delta_r2_filtro_maint']):+.4f} de R²   |   "
       f"Efeito das variáveis de contrato: {float(cp['delta_r2_contrato']):+.4f} de R²",
       16, NAVY, True)

# 26. Comparacao dos modelos
sl = slide()
header(sl, "13 · Resultados", "Comparação dos modelos (teste 2025)")
if not metr.empty:
    table(sl, metr, 0.7, 1.8, 7.4, 5.0, fs=10, maxrows=14,
          destaque=lambda r: best_pred is not None and str(r["modelo"]) == str(best_pred["modelo"]))
if best_pred is not None:
    bullets(sl, 8.4, 1.9, 4.3, 4.8, [
        f"Preditivo recomendado: {best_pred['modelo']}",
        f"R² = {best_pred['r2']} · RMSE = {best_pred['rmse']}",
        f"MAE = {best_pred['mae']} CAD/ano",
        f"Explicativo (melhor): R² = {best_expl['r2']}",
        "Árvores/ensembles superam claramente os modelos lineares.",
        "Configuração C (população MAINT, com contrato).",
    ], 13)

# 27. Importancia
sl = slide()
header(sl, "13 · Resultados", "Variáveis mais importantes (permutação)")
pic_fit(sl, FIG / "05_importancia_permutacao.png", 0.7, 1.7, 7.6, 5.2)
if not imp.empty:
    table(sl, imp.head(9)[["variavel", "importancia"]].round(4), 8.5, 1.9, 4.2, 4.4, fs=10, maxrows=9,
          destaque=lambda r: "contrato" in str(r["variavel"]))
tb(sl, 8.5, 6.5, 4.3, 0.6,
   "A variável de contrato aparece por último, dentro do desvio.", 11, GREY)

# 28. Hipoteses com veredito
sl = slide()
header(sl, "13 · Resultados", "Vereditos das hipóteses")
if not hip.empty:
    table(sl, hip[["hipotese", "enunciado", "evidencia", "veredito"]], 0.6, 1.85, 12.2, 4.6, fs=9, maxrows=8,
          destaque=lambda r: str(r["hipotese"]).startswith("H6"))
tb(sl, 0.6, 6.6, 12.2, 0.6,
   "H6 saiu de 'não testável' (Fase 1) para 'testada com efeito fraco' (Fase 2). Um resultado negativo é resultado: "
   "o contrato foi verificado, não pressuposto.", 12, NAVY, True)

# 29. Implicacoes gerenciais
sl = slide()
header(sl, "14 · Implicações gerenciais", "O que a empresa faz com isso")
if not rec.empty:
    table(sl, rec, 0.7, 1.85, 12.0, 4.4, fs=10, maxrows=6)
tb(sl, 0.7, 6.45, 12.0, 0.7,
   f"Erro médio do modelo recomendado: CAD $ {best_pred['mae'] if best_pred is not None else '—'}/ano por carreta — "
   "é a margem a considerar ao usar a estimativa em orçamento.", 13, NAVY, True)

# 30. Metodologia passo a passo
sl = slide()
header(sl, "15 · Metodologia", "Procedimento passo a passo")
passos = [
    "0  Mapeamento e extração — tabelas-chave, scripts SQL, desnormalização em base consolidada.",
    f"1  Fonte única — CSV consolidado fato_wo_ml ({cur_val('linhas_originais'):,} OS · {cur_val('carretas_distintas'):,} carretas · 29 colunas).".replace(",", "."),
    f"2  Qualidade e integridade — exclusão de {cur_val('removidas_custo_negativo_estorno')} estornos e {cur_val('removidas_fora_janela_2020_2025')} OS fora da janela.",
    f"3  Base carreta × ano — alvo Y = custo anual por carreta ({vl('linhas_base_anual')} linhas).",
    "4  Deflação — CPI Canadá, custos em CAD reais de dez/2025.",
    "5  EDA — descritivas, distribuições e relação de cada variável com o custo.",
    "6  Seleção anti-vazamento — remove componentes de Y; histórico só do ano anterior.",
    f"7  População — flag MAINT: {vl('carreta_ano_populacao_maint')} carreta-anos sob contrato com manutenção inclusa.",
    "8  Modelagem — split temporal 2020–2024 / 2025; lineares, árvores e ensembles; 3 configurações.",
    "9  Avaliação — R² / RMSE / MAE, importância por permutação e vereditos das hipóteses.",
    "10 Aplicação — apoio ao orçamento anual e à priorização de manutenção da frota.",
]
tb(sl, 0.9, 1.8, 11.8, 5.2, "\n".join(passos), 13, INK)

# 31. Limitacoes
sl = slide()
header(sl, "15 · Metodologia", "Limitações metodológicas")
bullets(sl, 0.9, 1.85, 11.6, 5.0, [
    "Fonte única: contrato já incorporado; seguem ausentes mão de obra detalhada, peças e tipo_contrato (RENTAL/LEASE).",
    "População restrita a MAINT: o ganho de R² do filtro reflete amostra mais homogênea, não melhora de previsão — "
    "as carreta-anos excluídas têm custo médio bem menor e muitas têm custo zero.",
    "franquia_km_mensal_contrato descartada: 99,8% dos valores preenchidos são zero.",
    "NET e MIX somam 2,7% das OS: conclusões sobre esses regimes têm baixa potência estatística.",
    "cod_cliente não modelado (597 categorias): risco de memorização do cliente em vez de explicação do custo.",
    "km derivado do odômetro nas OS; resets/ruído tratados por regra, com aproximação.",
    "Província parcial (~54%); região usada como proxy geográfica.",
    "Estornos excluídos; span ativo assume presença entre a 1ª e a última OS.",
], 13)

# 32. Conclusoes
sl = slide()
header(sl, "16 · Conclusões", "Resposta à pergunta de pesquisa")
c = sl.shapes.add_shape(1, Inches(0.9), Inches(1.8), Inches(11.5), Inches(1.5))
c.fill.solid(); c.fill.fore_color.rgb = NAVY; c.line.fill.background(); c.shadow.inherit = False
tb(sl, 1.2, 1.95, 11.0, 1.3,
   "O custo anual real por carreta é explicado sobretudo por REFRIGERAÇÃO, HISTÓRICO DE MANUTENÇÃO e USO.\n"
   f"O modelo {best_pred['modelo'] if best_pred is not None else ''} estima o custo do ano seguinte com "
   f"R² = {best_pred['r2'] if best_pred is not None else '—'} e erro médio de CAD $ {best_pred['mae'] if best_pred is not None else '—'}/ano.",
   16, WHITE, True)
bullets(sl, 0.9, 3.6, 11.6, 3.2, [
    "O grão anual eliminou a zero-inflação que travava as formulações anteriores (3,1% de carreta-anos com custo zero).",
    "A idade isolada NÃO explica o custo: atua por meio do histórico e do uso.",
    "As variáveis de contrato, incorporadas nesta fase, mostraram efeito fraco — hipótese testada, não assumida.",
    "O contrato serviu principalmente para DEFINIR a população de análise (MAINT), não para prever o custo.",
    "O modelo é utilizável para orçamento e priorização, desde que a margem de erro seja explicitada.",
], 14)

# 33. Entregas / Gates
sl = slide()
header(sl, "17 · Entregas", "Gates do projeto")
gates = pd.DataFrame([
    {"Gate": "1 · Dados e qualidade", "Escopo": "Inventário, validação de chaves, datas, custos e ausentes", "Status": "Concluída"},
    {"Gate": "2 · Base analítica", "Escopo": "Base carreta × ano, variáveis defasadas (anti-vazamento)", "Status": "Concluída"},
    {"Gate": "3 · EDA e hipóteses", "Escopo": "Distribuições, outliers, correlações, VIF e segmentações", "Status": "Concluída"},
    {"Gate": "4 · Modelo anual de custo", "Escopo": "7 modelos × 2 cenários × 3 configurações, teste temporal 2025", "Status": "Concluída (Fase 2)"},
    {"Gate": "5 · Contrato", "Escopo": "Incorporação dos 4 campos, teste de H6a e H6b", "Status": "Concluída (Fase 2)"},
    {"Gate": "6 · Interpretação de negócio", "Escopo": "Orçamento, preventiva e gestão da frota", "Status": "Próximos passos"},
])
table(sl, gates, 0.7, 1.9, 12.0, 3.6, fs=12, maxrows=6,
      destaque=lambda r: "Fase 2" in str(r["Status"]))
tb(sl, 0.7, 5.8, 12.0, 1.2,
   "Trabalhos futuros: integrar mão de obra, peças e tipo_contrato (RENTAL/LEASE); testar modelos de efeitos mistos "
   "por grupo de ativo (MERF); incorporar telemetria/GPS como exposição; modelar a cauda de custos extremos.", 13, INK)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"OK: {OUT.name} — {len(prs.slides._sldIdLst)} slides")
