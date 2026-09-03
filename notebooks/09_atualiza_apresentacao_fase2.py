"""09 - Consolidacao da apresentacao da Fase 2.

Parte de `docs/entregas/Apresentacao_QuatroNorte_agosto.pptx` (o PowerPoint de origem do
PDF apresentado em 2026-08-05), atualiza os numeros com os resultados da Fase 2,
reconstroi as tabelas a partir de `reports/`, reexibe os slides de modelagem que estavam
ocultos e acrescenta o bloco de contrato.

Revisao de 2026-09-02: o bloco da Fase 2 deixou de narrar as tres configuracoes
(sem contrato -> com contrato -> recorte MAINT) e o experimento de divisao por
refrigeracao. Passa a apresentar os dois alvos do projeto, Y1 (custo anual) e Y2
(numero de OS no ano). Os 34 slides da entrega 1 seguem intactos e na ordem original.

Saida: `docs/entregas/Apresentacao_QuatroNorte_Fase2.pptx`.
O arquivo de origem NAO e modificado.
"""
import copy
import pandas as pd
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
TABLES = ROOT / "reports" / "tables"
FIG = ROOT / "reports" / "figures"
FIG_EDA = FIG / "eda"
SRC = ROOT / "docs" / "entregas" / "Apresentacao_QuatroNorte_agosto.pptx"
OUT = ROOT / "docs" / "entregas" / "Apresentacao_QuatroNorte_Fase2.pptx"


def rd(n):
    p = TABLES / n
    return pd.read_csv(p) if p.exists() else pd.DataFrame()


stats = rd("03c_stats_ppt.csv").set_index("metrica")["valor"]
metr = rd("05_metricas_modelos.csv")
cmpcfg = rd("05_comparacao_configuracoes.csv")
best = rd("05_modelo_recomendado.csv")
imp = rd("05_importancia_permutacao_random_forest.csv")
hip = rd("06_hipoteses_final.csv")
sel = rd("05_selecao_variaveis.csv")
corr = rd("03b_correlacao_com_y.csv")
eta = rd("03b_eta_categoricas.csv")
vif = rd("03b_vif.csv")
est = rd("03b_estatisticas_descritivas.csv")
val = rd("02_validacao_base_anual.csv")
rec = rd("06_recomendacoes_negocio.csv")
metr_all = rd("05_metricas_por_configuracao.csv")
# alvos decompostos (D8, 2026-09-02)
alvos_dec = rd("05_comparacao_alvos_decompostos.csv")
alvos_y23 = rd("05_metricas_alvos_y2_y3.csv")
# curadoria de features (02/09)
cur = rd("12_curadoria_decisoes.csv")
cur_cmp = rd("12_comparativo_Y1_Y2.csv")
cur_dum = rd("12_mapa_dummies.csv")


def sv(k):
    v = stats[k]
    try:
        return float(v)
    except (TypeError, ValueError):
        return v


def vl(chave, d="—"):
    r = val.loc[val["checagem"] == chave, "valor"]
    return r.iloc[0] if len(r) else d


best_pred = best[best["cenario"] == "preditivo"].iloc[0]
best_expl = best[best["cenario"] == "explicativo"].iloc[0]
cp = cmpcfg[cmpcfg["cenario"] == "preditivo"].iloc[0]
ce = cmpcfg[cmpcfg["cenario"] == "explicativo"].iloc[0]

# --- alvos Y1 e Y2 (D8) -------------------------------------------------------
_dec = alvos_dec[alvos_dec["caminho"].str.startswith("decomposto")].iloc[0]
_dir = alvos_dec[alvos_dec["caminho"].str.startswith("direto")].iloc[0]
_r2_dec = float(_dec["r2"])
_y2 = alvos_y23[alvos_y23["alvo"] == "n_os_ano"].sort_values("rmse").iloc[0]
_y3 = alvos_y23[alvos_y23["alvo"] == "custo_medio_por_os_ano"].sort_values("rmse").iloc[0]
_ref = imp[imp["variavel"] == "flag_refrigerado"]
_imp_ref = float(_ref.iloc[0]["importancia"]) if len(_ref) else 0.21


def br(v, casas=4):
    """Formata numero com virgula decimal, como manda o portugues do deck."""
    return f"{float(v):.{casas}f}".replace(".", ",")

NAVY = RGBColor(0x14, 0x2B, 0x45)
INK = RGBColor(0x1C, 0x24, 0x33)
ORANGE = RGBColor(0xE8, 0x71, 0x3A)
PAPER = RGBColor(0xF6, 0xF1, 0xE7)
GREY = RGBColor(0x5B, 0x63, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(str(SRC))
slides = list(prs.slides)
print(f"origem: {len(slides)} slides")

# ============================================================
# 1) Substituicoes de texto (numeros e afirmacoes superadas)
# ============================================================
SUBS = [
    ("223.590", "217.217"),
    ("9.859", "9.585"),
    ("CAD $ 82.428 M", f"CAD $ {sv('custo_total_real_mi')} M"),
    ("82.428", str(sv("custo_total_real_mi"))),
    ("77.18 mi nominal", f"{sv('custo_total_nominal_mi')} mi nominal"),
    ("47.666", f"{int(vl('linhas_base_anual')):,}".replace(",", ".")),
    ("1673.72", f"{sv('y_media'):.2f}"),
    ("812.55", f"{sv('y_mediana'):.2f}"),
    ("4221.76", f"{sv('y_p90'):.2f}"),
    ("Assimetria 3.79", f"Assimetria {sv('assimetria')}"),
    ("3.23% de carreta", f"{sv('share_y_zero_pct')}% de carreta"),
    # afirmacoes que se inverteram
    ("Base consolidada inicial", "Base consolidada única (29 colunas)"),
    ("Fonte inicial do estudo", "Fonte única do estudo"),
    ("Fonte inicial:", "Fonte única:"),
    ("Seleção inicial das variáveis do modelo", "Seleção das variáveis do modelo"),
    ("sugestão inicial", "decisão"),
    ("Comparação inicial dos modelos", "Comparação dos modelos"),
    ("São 25 variáveis explicativas candidatas da fonte inicial.",
     "São 29 colunas na fonte única, já incluindo o bloco de contrato."),
    ("Veredito Preliminar", "Veredito"),
    ("sem dados de contrato, mão de obra detalhada e peças — reduz o conjunto explicativo",
     "contrato incorporado na Fase 2; seguem ausentes mão de obra, peças e tipo_contrato (RENTAL/LEASE)"),
    ("Sem filtro por tipo de manutenção: analisa-se todo o custo interno.",
     "População: a frota completa. Restringir ao regime MAINT foi testado e alterou o R² em apenas +0,0017 — "
     "amostra mais homogênea, não previsão melhor."),
    ("Extração SQL, modelo estrela, joins e feature engineering .",
     "Extração SQL, modelo estrela, joins e feature engineering (etapa anterior de preparação)."),
    ("*Fase 2", "✓ concluído na Fase 2"),
]

n_sub = 0
for sl in slides:
    for sh in sl.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                novo = run.text
                for a, b in SUBS:
                    if a in novo:
                        novo = novo.replace(a, b)
                if novo != run.text:
                    run.text = novo
                    n_sub += 1
print(f"substituicoes de texto aplicadas em {n_sub} runs")


# ============================================================
# 2) Reconstrucao de tabelas a partir de reports/
# ============================================================
def set_cell(cell, texto, bold=None):
    tf = cell.text_frame
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = str(texto)
        for r in p0.runs[1:]:
            r.text = ""
        if bold is not None:
            p0.runs[0].font.bold = bold
    else:
        cell.text = str(texto)


def clonar_linha(tbl, idx_modelo=1):
    """Duplica uma linha existente preservando formatacao."""
    tr = copy.deepcopy(tbl._tbl.tr_lst[idx_modelo])
    tbl._tbl.append(tr)


def preencher(tbl, df, cabecalho=True, ajustar=True):
    """Preenche a tabela com o df, ajustando o numero de linhas."""
    n_alvo = len(df) + (1 if cabecalho else 0)
    if ajustar:
        while len(tbl.rows) < n_alvo:
            clonar_linha(tbl)
        while len(tbl.rows) > n_alvo:
            tbl._tbl.remove(tbl._tbl.tr_lst[-1])
    ncols = len(tbl.columns)
    if cabecalho:
        for j, col in enumerate(df.columns[:ncols]):
            set_cell(tbl.cell(0, j), acentuar(col))
    off = 1 if cabecalho else 0
    for i in range(len(df)):
        for j in range(min(ncols, df.shape[1])):
            set_cell(tbl.cell(i + off, j), acentuar(df.iloc[i, j]))


def tabelas(sl):
    return [sh.table for sh in sl.shapes if sh.has_table]


ACENTOS = {
    "removida": "removida", "mantida": "mantida",
    "mantida (so explicativo)": "mantida (só explicativo)",
    "removida do modelo": "removida do modelo",
    "explicativo / inicio_ano no preditivo": "explicativo / início de ano no preditivo",
    "condicional": "condicional",
    "variavel": "variável", "decisao": "decisão", "cenario": "cenário",
    "importancia": "importância", "hipotese": "hipótese", "enunciado": "enunciado",
    "veredito": "veredito", "papel": "papel", "recomendacao": "recomendação",
    "desvio_padrao": "desvio padrão",
}


def acentuar(v):
    s_ = str(v)
    return ACENTOS.get(s_, s_)


# -- slide 5: hipoteses (agora com veredito final, H6a/H6b)
t = tabelas(slides[4])[0]
hip_v = hip[["hipotese", "enunciado", "veredito"]].copy()
hip_v.columns = ["Hipótese", "Descrição", "Veredito"]
preencher(t, hip_v)
print("slide 5: hipoteses ->", len(hip_v), "linhas")

# -- slide 11: estatisticas descritivas
t = tabelas(slides[10])[0]
cols = [c for c in ["variavel", "N", "media", "mediana", "desvio_padrao", "min", "max", "assimetria"] if c in est.columns]
ncol = len(tabelas(slides[10])[0].columns)
preencher(t, est[cols].head(16).iloc[:, :ncol])
print("slide 11: descritivas atualizadas")

# -- slide 20: correlacao e eta
tt = tabelas(slides[19])
cc = corr[["variavel", "spearman", "papel"]].sort_values("spearman", key=abs, ascending=False).head(15)
cc.columns = ["variável", "spearman", "papel"]
preencher(tt[0], cc.iloc[:, :len(tt[0].columns)])
ee = eta[["variavel", "eta", "n_categorias"]].copy()
ee.columns = ["variável", "eta", "n_categorias"]
preencher(tt[1], ee.iloc[:, :len(tt[1].columns)])
print("slide 20: correlacao/eta atualizados")

# -- slide 22: VIF
t = [x for x in tabelas(slides[21])][0]
vv = vif.copy()
vv.columns = ["variável", "vif"][: len(vv.columns)]
preencher(t, vv.iloc[:, :len(t.columns)])
print("slide 22: VIF atualizado")

# -- slide 23: selecao de variaveis (duas colunas)
tt = tabelas(slides[22])
ss = sel[["variavel", "decisao"]].copy()
ss.columns = ["variável", "decisão"]
meta = (len(ss) + 1) // 2
preencher(tt[0], ss.iloc[:meta])
preencher(tt[1], ss.iloc[meta:])
print("slide 23: selecao ->", len(ss), "variaveis")

# -- slide 27: metricas dos modelos
t = tabelas(slides[26])[0]
mm = metr.copy()
mm.columns = ["cenário", "modelo", "r2", "rmse", "mae"]
preencher(t, mm)
for sh in slides[26].shapes:
    if sh.has_text_frame and "Preditivo recomendado" in sh.text_frame.text:
        tf = sh.text_frame
        novos = [
            f"Preditivo recomendado: {best_pred['modelo']}",
            f"R² = {best_pred['r2']} · RMSE = {best_pred['rmse']}",
            f"MAE = {best_pred['mae']} CAD/ano",
            f"Explicativo (melhor): R² = {best_expl['r2']}",
            "Árvores/ensembles superam modelos lineares.",
            f"Efeito das variáveis de contrato: +{br(cp['delta_r2_contrato'])} de R².",
            f"Caminho decomposto (nº de OS × custo médio): R² = {br(_r2_dec)}.",
        ]
        base_p = tf.paragraphs[0]
        while len(tf.paragraphs) > 1:
            tf._txBody.remove(tf.paragraphs[-1]._p)
        for i, txt in enumerate(novos):
            p = base_p if i == 0 else tf.add_paragraph()
            if p.runs:
                p.runs[0].text = "•  " + txt
                for r in p.runs[1:]:
                    r.text = ""
            else:
                r = p.add_run()
                r.text = "•  " + txt
                r.font.size = Pt(13)
                r.font.color.rgb = INK
                r.font.name = "Calibri"
print("slide 27: metricas atualizadas")

# -- slide 28: importancia por permutacao
t = tabelas(slides[27])[0]
ii = imp.head(9)[["variavel", "importancia"]].round(4).copy()
ii.columns = ["variável", "importância"]
preencher(t, ii)
for sh in slides[27].shapes:
    if sh.shape_type == 13 or sh.__class__.__name__ == "Picture":
        el = sh._element
        el.getparent().remove(el)
slides[27].shapes.add_picture(str(FIG / "05_importancia_permutacao.png"),
                              Inches(0.7), Inches(1.7), Inches(7.6), Inches(5.2))
print("slide 28: importancias atualizadas")

# (gravação no final do arquivo)


# ============================================================
# 3) Reexibir os slides de modelagem que estavam ocultos
# ============================================================
REEXIBIR = [26, 27, 28]  # 27 metricas, 28 importancias, 29 "O que dizem os dados"
for i in REEXIBIR:
    el = slides[i]._element
    if el.get("show") == "0":
        el.set("show", "1")
print("slides reexibidos:", [i + 1 for i in REEXIBIR])

# slide 29 - "O que dizem os dados": atualizar a leitura frente a literatura
for sh in slides[28].shapes:
    if not sh.has_text_frame:
        continue
    for para in sh.text_frame.paragraphs:
        for run in para.runs:
            t0 = run.text
            t1 = (t0.replace("R² 0,43 (preditivo) a 0,57 (explicativo)",
                             f"R² {float(best_pred['r2']):.2f} (preditivo) a {float(best_expl['r2']):.2f} (explicativo)")
                    .replace("regressões lineares ficam com R²<0",
                             "regressões lineares ficam muito atrás (R² ~ 0,27)")
                    .replace("idade isolada ρ=0,02", "idade isolada ρ=0,03")
                    .replace("removê-la derruba o R² em 0,22", "removê-la derruba o R² em 0,17")
                    .replace("o R² cai de 0,57 → 0,43",
                             f"o R² cai de {float(best_expl['r2']):.2f} para {float(best_pred['r2']):.2f}"))
            if t1 != t0:
                run.text = t1

# ============================================================
# 4) Novos slides da Fase 2 (mesma identidade visual)
# ============================================================
LAYOUT = next((l for l in prs.slide_layouts if l.name == "Blank"), prs.slide_layouts[6])


def novo_slide(bg=PAPER):
    sl = prs.slides.add_slide(LAYOUT)
    r = sl.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    sl.shapes._spTree.remove(r._element)
    sl.shapes._spTree.insert(2, r._element)
    return sl


def txt(sl, x, y, w, h, texto, size=16, cor=INK, bold=False):
    box = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, linha in enumerate(str(texto).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = linha
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = cor
        r.font.name = "Calibri"
    return box


def cabecalho(sl, kicker, titulo):
    bar = sl.shapes.add_shape(1, Inches(0.6), Inches(0.55), Inches(0.16), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()
    bar.shadow.inherit = False
    txt(sl, 0.9, 0.5, 11.8, 0.4, kicker.upper(), 12, ORANGE, True)
    txt(sl, 0.9, 0.86, 11.8, 0.9, titulo, 27, NAVY, True)


def marcadores(sl, x, y, w, h, itens, size=15):
    box = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, it in enumerate(itens):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        r = p.add_run()
        r.text = "•  " + it
        r.font.size = Pt(size)
        r.font.color.rgb = INK
        r.font.name = "Calibri"
    return box


def nova_tabela(sl, df, x, y, w, h, fs=11):
    nr, nc = df.shape
    gt = sl.shapes.add_table(nr + 1, nc, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for j, col in enumerate(df.columns):
        c = gt.cell(0, j)
        c.text = str(col)
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
        p = c.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(fs)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE
        p.runs[0].font.name = "Calibri"
    for i in range(nr):
        for j in range(nc):
            c = gt.cell(i + 1, j)
            c.text = str(df.iloc[i, j])
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(0xEE, 0xEA, 0xDE)
            p = c.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].font.size = Pt(fs)
                p.runs[0].font.color.rgb = INK
                p.runs[0].font.name = "Calibri"
    return gt


def kpi(sl, x, y, w, h, valor, rotulo, cor=ORANGE):
    c = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid()
    c.fill.fore_color.rgb = WHITE
    c.line.color.rgb = NAVY
    c.shadow.inherit = False
    txt(sl, x + 0.2, y + 0.08, w - 0.4, 0.45, valor, 22, cor, True)
    txt(sl, x + 0.2, y + 0.58, w - 0.4, 0.30, rotulo, 11, GREY)


novos = []

# A. Os quatro campos de contrato
sl = novo_slide()
novos.append(("contrato_campos", sl))
cabecalho(sl, "Fase 2 · dados novos", "Contrato: os quatro campos que chegaram à base")
txt(sl, 0.9, 1.68, 11.6, 0.4,
    "H6 estava declarada desde a fase anterior, sem dados para testá-la. Os campos de contrato passaram a integrar a fonte única.",
    12, GREY)
campos = pd.DataFrame([
    {"campo": "tempo_contrato_meses_ate_reparo", "cobertura": "92,5%",
     "perfil": "mediana 35,1 meses · sem valores negativos", "destino": "MODELADA — responde H6a"},
    {"campo": "tipo_manutencao", "cobertura": "92,5%",
     "perfil": "MAINT 89,7% · MIX 1,6% · NET 1,1%", "destino": "variável do modelo (dummy); testa H6b"},
    {"campo": "cod_cliente", "cobertura": "77,2%",
     "perfil": "597 clientes distintos", "destino": "descritiva — fora do modelo"},
    {"campo": "franquia_km_mensal_contrato", "cobertura": "71,6%",
     "perfil": "99,8% dos preenchidos são ZERO", "destino": "DESCARTADA — variância nula"},
])
nova_tabela(sl, campos, 0.7, 2.2, 12.0, 2.2, fs=11)
marcadores(sl, 0.9, 4.8, 11.6, 2.1, [
    "A fonte permanece ÚNICA: os campos vieram dentro do próprio CSV consolidado, sem join novo.",
    "Contrato não é atributo fixo do ativo — 51,5% das carretas mudam de regime no período, o que exigiu agregação por carreta-ano.",
    "O regime de manutenção entrou como VARIÁVEL do modelo, e não como recorte da amostra: a frota completa é a população.",
], 14)

# B. EDA de contrato
sl = novo_slide()
novos.append(("contrato_eda", sl))
cabecalho(sl, "Fase 2 · EDA de contrato", "Distribuições das variáveis de contrato")
for nome, fx, fy in [("quant_tempo_contrato_meses_fim_ano.png", 0.8, 1.8),
                     ("quant_share_maint_ano.png", 6.9, 1.8),
                     ("quant_n_clientes_ano.png", 0.8, 4.4),
                     ("quant_trocou_contrato_ano.png", 6.9, 4.4)]:
    p = FIG_EDA / nome
    if p.exists():
        sl.shapes.add_picture(str(p), Inches(fx), Inches(fy), height=Inches(2.4))

# C. Boxplot por regime contratual (H6b)
sl = novo_slide()
novos.append(("contrato_h6b", sl))
cabecalho(sl, "Fase 2 · EDA de contrato", "Custo anual por regime contratual — evidência de H6b")
p = FIG_EDA / "quali_tipo_manutencao_ano.png"
if p.exists():
    sl.shapes.add_picture(str(p), Inches(0.9), Inches(1.85), height=Inches(4.9))
_et = eta.loc[eta["variavel"] == "tipo_manutencao_ano", "eta"]
_et = float(_et.iloc[0]) if len(_et) else float("nan")
marcadores(sl, 8.6, 2.0, 4.2, 4.6, [
    f"eta = {_et:.3f}: separação fraca entre regimes.",
    "MAINT concentra 89,7% das OS; NET e MIX somam 2,7% — baixa potência estatística.",
    "O tipo de contrato deixou de recortar a amostra e passou a ser VARIÁVEL do modelo — é assim que H6b é testada.",
    "Diferença entre regimes pode refletir quem paga o reparo, não quanto ele custa.",
], 13)

# C2. Curadoria de features (02/09)
sl = novo_slide()
novos.append(("curadoria", sl))
cabecalho(sl, "Fase 2 · curadoria de variáveis", "As 40 variáveis revisadas uma a uma (02/set)")
txt(sl, 0.9, 1.66, 11.6, 0.42,
    "O Grupo revisou cada variável desenhada na fase anterior — manter, remover ou transformar em dummy — e a EDA foi refeita "
    "para confrontar cada decisão com a evidência.", 12, GREY)

_n_modelo = int((cur["entra_no_modelo"] == "sim").sum()) if not cur.empty else 0
_n_fora = int((cur["entra_no_modelo"] != "sim").sum()) if not cur.empty else 0
_n_dummy = len(cur_dum) if not cur_dum.empty else 0
_n_dummy_var = int(cur_dum["variavel_origem"].nunique()) if not cur_dum.empty else 0
kpi(sl, 0.9, 2.25, 3.5, 1.05, str(_n_modelo), "variáveis no modelo", ORANGE)
kpi(sl, 4.8, 2.25, 3.5, 1.05, str(_n_fora), "retiradas ou promovidas a alvo", NAVY)
kpi(sl, 8.7, 2.25, 3.5, 1.05, str(_n_dummy), f"colunas dummy, de {_n_dummy_var} variáveis", NAVY)

curad = pd.DataFrame([
    {"decisão": "DUMMY (one-hot)",
     "variáveis": "descricao_carreta · unit_subtype · flag_refrigerado · tipo_manutencao_ano · vmrs_predominante_ano",
     "evidência": "η de 0,18 a 0,59"},
    {"decisão": "MANTER — histórico defasado",
     "variáveis": "nº de OS e custo do ano anterior e acumulados · anos ativos",
     "evidência": "ρ de 0,26 a 0,54"},
    {"decisão": "REMOVER — geografia (H5 encerrada)",
     "variáveis": "regiao_operacao · provincia_estado",
     "evidência": "η ≤ 0,14"},
    {"decisão": "REMOVER — atributos de efeito fraco",
     "variáveis": "cod_montadora · suspension_type · new_used_indicator · ano_modelo · comprimento · share_pm_ano",
     "evidência": "η e ρ ≤ 0,23"},
    {"decisão": "REMOVER — risco de memorização",
     "variáveis": "cod_cliente_predominante_ano (597 categorias)",
     "evidência": "η 0,61, mas não generaliza"},
    {"decisão": "PROMOVER A ALVO",
     "variáveis": "n_os_ano (Y2) · custo_medio_por_os_ano",
     "evidência": "são componentes de Y1"},
])
nova_tabela(sl, curad, 0.7, 3.55, 12.0, 2.35, fs=10)
marcadores(sl, 0.9, 6.02, 11.6, 1.0, [
    "A premissa do Grupo para VMRS — média dos últimos 5 anos — foi testada contra três alternativas e VENCEU (ρ 0,474).",
    "Cinco variáveis fracas em Y1 mostraram-se relevantes em Y2: decisão em aberto para a próxima rodada.",
], 12)

# D. Os dois alvos: Y1 e Y2
sl = novo_slide()
novos.append(("alvos", sl))
cabecalho(sl, "Fase 2 · desenho da modelagem", "Y1 e Y2: o que o projeto prevê")
txt(sl, 0.9, 1.68, 11.6, 0.4,
    "O custo anual é o produto de duas quantidades: quantas ordens de serviço a carreta gera e quanto custa cada uma. Prevemos as duas.",
    12, GREY)
alvos_tab = pd.DataFrame([
    {"alvo": "Y1 · custo_ano_real", "o que prevê": "custo anual de manutenção por carreta (CAD real)",
     "para que serve": "base do orçamento e da provisão da frota"},
    {"alvo": "Y2 · n_os_ano", "o que prevê": "quantidade de ordens de serviço no ano",
     "para que serve": "capacidade de oficina; é o componente que mais move o custo"},
])
nova_tabela(sl, alvos_tab, 0.9, 2.25, 11.6, 1.35, fs=12)
kpi(sl, 0.9, 4.0, 3.5, 1.15, br(_dir["r2"]), "Y1 · direto", NAVY)
kpi(sl, 4.8, 4.0, 3.5, 1.15, br(_r2_dec), "Y1 · decomposto", ORANGE)
kpi(sl, 8.7, 4.0, 3.5, 1.15, br(_y2["r2"]), "Y2 · nº de OS", ORANGE)
txt(sl, 0.9, 5.45, 11.6, 0.5,
    "Prever as partes e multiplicar supera prever o total direto: "
    f"{br(_dir['r2'])} → {br(_r2_dec)} de R².", 16, NAVY, True)
txt(sl, 0.9, 6.05, 11.6, 0.9,
    "Y1 = nº de OS × custo médio por OS é identidade EXATA na base (diferença máxima de 0,00 CAD em 47.715 linhas).\n"
    "Por isso os componentes não entram como variáveis explicativas de Y1 — usá-los seria a conta de volta, com R² = 1,0. "
    "Eles são previstos a partir do histórico e multiplicados.", 12, GREY)

# E. Implicacoes gerenciais
sl = novo_slide()
novos.append(("implicacoes", sl))
cabecalho(sl, "Item 13 · implicações gerenciais", "O que a empresa faz com estes resultados")
if not rec.empty:
    nova_tabela(sl, rec, 0.7, 1.9, 12.0, 4.2, fs=10)
txt(sl, 0.7, 6.35, 12.0, 0.6,
    f"Erro médio do modelo recomendado (Y1 decomposto): CAD $ {float(_dec['mae']):.0f}/ano por carreta — margem a declarar ao usar a estimativa em orçamento.",
    13, NAVY, True)

# F. Conclusoes
sl = novo_slide()
novos.append(("conclusoes", sl))
cabecalho(sl, "Item 16 · conclusões", "Resposta à pergunta de pesquisa")
c = sl.shapes.add_shape(1, Inches(0.9), Inches(1.8), Inches(11.5), Inches(1.45))
c.fill.solid()
c.fill.fore_color.rgb = NAVY
c.line.fill.background()
c.shadow.inherit = False
txt(sl, 1.2, 1.95, 11.0, 1.25,
    "O custo anual real por carreta é explicado sobretudo por REFRIGERAÇÃO, HISTÓRICO DE MANUTENÇÃO e USO.\n"
    f"Prevemos Y1 (custo anual) com R² = {br(_r2_dec)} e erro médio de CAD $ {float(_dec['mae']):.0f}/ano, "
    f"e Y2 (nº de OS no ano) com R² = {br(_y2['r2'])}.", 16, WHITE, True)
marcadores(sl, 0.9, 3.5, 11.6, 3.3, [
    f"O grão anual eliminou a zero-inflação que travava as formulações anteriores ({sv('share_y_zero_pct')}% de carreta-anos com custo zero).",
    "A idade isolada NÃO explica o custo: atua por meio do histórico e do uso.",
    f"Prever as partes e multiplicar supera prever o total: R² {br(_dir['r2'])} → {br(_r2_dec)}.",
    f"O nº de OS por ano é previsível (R² {br(_y2['r2'], 3)}); o custo médio por OS é o elo fraco (R² {br(_y3['r2'], 3)}).",
    "Do contrato, o que importa é o TIPO (4ª variável mais importante), não a duração — cuja contribuição é nula.",
    "O modelo é utilizável para orçamento e priorização, desde que a margem de erro seja declarada.",
], 14)


# ---- G. Linha do tempo: o que mudou desde a apresentação ----
sl = novo_slide()
novos.append(("linha_do_tempo", sl))
cabecalho(sl, "Fase 2 · o que mudou", "Linha do tempo desde a apresentação de agosto")
txt(sl, 0.9, 1.68, 11.6, 0.4,
    "A Fase 1 encerrou na seleção de variáveis, com a modelagem marcada como Fase 2. Um dado novo mudou o que era possível testar.",
    12, GREY)
linha = pd.DataFrame([
    {"quando": "05/ago", "o que aconteceu": "Apresentação da Fase 1 (28 slides)",
     "situação": "H6 declarada, sem dados; Gate 4 em aberto"},
    {"quando": "16/ago", "o que aconteceu": "Base reextraída, com dados de contrato",
     "situação": "25 → 29 colunas · 223.590 → 217.217 OS · 9.859 → 9.585 carretas"},
    {"quando": "02/set", "o que aconteceu": "Curadoria de variáveis pelo Grupo",
     "situação": "40 variáveis revisadas uma a uma; geografia encerrada (H5)"},
    {"quando": "02/set", "o que aconteceu": "População: a frota completa",
     "situação": "47.715 carreta-anos; o tipo de contrato passa a ser variável do modelo"},
    {"quando": "02/set", "o que aconteceu": "Alvo decomposto: Y1 e Y2",
     "situação": f"prever nº de OS e custo médio supera prever o custo direto ({br(_dir['r2'])} → {br(_r2_dec)})"},
])
nova_tabela(sl, linha, 0.7, 2.2, 12.0, 2.9, fs=10)
marcadores(sl, 0.9, 5.3, 11.6, 1.6, [
    "A fonte permanece ÚNICA: o contrato veio dentro do próprio CSV consolidado, sem join novo.",
    "Todos os slides da apresentação anterior foram mantidos, na mesma ordem — os números é que foram atualizados.",
], 14)

# ---- K. Resultados finais: Y1 e Y2 ----
sl = novo_slide()
novos.append(("finais", sl))
cabecalho(sl, "Item 12 · resultados finais", "Y1 e Y2 no teste de 2025")
txt(sl, 0.9, 1.62, 11.6, 0.45,
    "Treino 2020–2024, teste 2025. Só entram variáveis conhecidas no início do ano: atributos da carreta e histórico até o ano anterior.",
    12, GREY)
fin = pd.DataFrame([
    {"alvo / caminho": "Y1 · custo anual — direto", "modelo": str(_dir["modelo"]),
     "R²": br(_dir["r2"]), "RMSE": br(_dir["rmse"], 1), "MAE": br(_dir["mae"], 1)},
    {"alvo / caminho": "Y1 · custo anual — decomposto", "modelo": "nº de OS × custo médio",
     "R²": br(_r2_dec), "RMSE": br(_dec["rmse"], 1), "MAE": br(_dec["mae"], 1)},
    {"alvo / caminho": "Y2 · nº de OS no ano", "modelo": str(_y2["modelo"]),
     "R²": br(_y2["r2"]), "RMSE": br(_y2["rmse"], 2), "MAE": br(_y2["mae"], 2)},
])
nova_tabela(sl, fin, 0.7, 2.2, 7.6, 1.9, fs=11)
marcadores(sl, 8.5, 2.25, 4.3, 2.9, [
    f"Y2 é mais previsível que Y1: {br(_y2['r2'], 3)} contra {br(_dir['r2'], 3)}.",
    "Em Y2 a regressão linear chega a R² 0,577 — a 0,031 do Random Forest.",
    "Nº de OS é contagem com superdispersão (variância/média = 4,0).",
], 12)
c = sl.shapes.add_shape(1, Inches(0.7), Inches(5.35), Inches(12.0), Inches(1.65))
c.fill.solid()
c.fill.fore_color.rgb = NAVY
c.line.fill.background()
c.shadow.inherit = False
txt(sl, 1.0, 5.5, 11.4, 1.45,
    f"DOIS MODELOS ENTREGUES — Y1: custo anual por carreta, R² = {br(_r2_dec)} pelo caminho decomposto "
    f"(MAE {float(_dec['mae']):.0f} CAD/ano).   Y2: nº de OS no ano, R² = {br(_y2['r2'])}.\n"
    f"O caminho decomposto vence porque o nº de OS é bem previsto. O custo médio por OS é o elo fraco "
    f"(R² {br(_y3['r2'], 3)}) — e a próxima frente de trabalho.", 13, WHITE, True)


# ============================================================
# 5) Reordenar os novos slides para as posicoes certas
# ============================================================
sldIdLst = prs.slides._sldIdLst


def id_do_slide(sl):
    for x in list(sldIdLst):
        if prs.part.related_part(x.rId) is sl.part:
            return x
    raise ValueError("slide nao encontrado no sldIdLst")


# ---- L. Item 14 — Limitações atualizadas na Fase 2 ----
sl = novo_slide()
novos.append(("limitacoes_f2", sl))
cabecalho(sl, "Item 14 · limitações", "O que mudou nas limitações após a Fase 2")
txt(sl, 0.9, 1.62, 11.6, 0.45,
    "As limitações declaradas em agosto continuam válidas, exceto as que esta fase resolveu — e há três novas, criadas pelas próprias escolhas da Fase 2.",
    12, GREY)
lim = pd.DataFrame([
    {"limitação declarada em agosto": "Sem dados de contrato — reduz o conjunto explicativo",
     "situação após a Fase 2": "RESOLVIDA — contrato incorporado e H6 testada"},
    {"limitação declarada em agosto": "Sem filtro por tipo de manutenção",
     "situação após a Fase 2": "RESOLVIDA de outro modo — o tipo de manutenção entrou como VARIÁVEL do modelo, em vez de recortar a amostra"},
    {"limitação declarada em agosto": "km derivado do odômetro; província parcial (~54%)",
     "situação após a Fase 2": "permanece"},
])
nova_tabela(sl, lim, 0.7, 2.2, 12.0, 1.7, fs=11)
txt(sl, 0.9, 4.15, 11.6, 0.4, "Limitações novas, decorrentes das escolhas desta fase:", 14, NAVY, True)
marcadores(sl, 0.9, 4.6, 11.6, 2.3, [
    f"O custo médio por OS é pouco previsível (R² {br(_y3['r2'], 3)}): é o fator que limita a precisão do caminho decomposto.",
    "NET e MIX somam 2,7% das OS — conclusões sobre esses regimes têm baixa potência estatística.",
    "cod_cliente não foi modelado (597 categorias) e franquia_km_mensal_contrato foi descartada (99,8% de zeros).",
    "A avaliação usa um ano de teste por vez; a variação entre anos é maior que ganhos medidos dentro de um único ano.",
    "Nenhuma previsão de 2026 foi gerada ainda: os números são validação sobre 2025.",
], 13)

# ---- M. Item 15 — Recomendações para projetos futuros ----
sl = novo_slide()
novos.append(("futuros", sl))
cabecalho(sl, "Item 15 · projetos futuros", "Recomendações para as próximas etapas")
txt(sl, 0.9, 1.62, 11.6, 0.45,
    "Contrato sai desta lista: deixou de ser trabalho futuro e virou resultado.",
    12, GREY)
fut = pd.DataFrame([
    {"frente": "Dados", "recomendação": "Integrar mão de obra e peças para decompor o custo por origem; incorporar tipo_contrato (RENTAL/LEASE), único campo de contrato ainda ausente"},
    {"frente": "Alvo Y3", "recomendação": "Melhorar a previsão do custo médio por OS (R² 0,085) — é o gargalo do caminho decomposto e onde há maior folga de ganho"},
    {"frente": "Modelagem", "recomendação": "Em Y2, adotar Binomial Negativa: é contagem com superdispersão (variância/média = 4,0), e os coeficientes viram razão de taxa, interpretáveis pelo negócio"},
    {"frente": "Previsão 2026", "recomendação": "Projetar quilometragem e montar a linha de cada carreta para 2026, aplicando o caminho decomposto — etapa ainda não executada"},
    {"frente": "Exposição", "recomendação": "Incorporar quilometragem planejada e telemetria (GPS) como medida de uso"},
])
nova_tabela(sl, fut, 0.7, 2.2, 12.0, 3.6, fs=11)
txt(sl, 0.7, 6.1, 12.0, 0.7,
    "A prioridade é a decomposição do custo por origem (mão de obra × peças): é o que permitiria dizer não só quanto custa, mas por quê.",
    13, NAVY, True)

# Os 34 slides apresentados em agosto permanecem CONTIGUOS e na ordem original.
# Todos os slides da Fase 2 entram DEPOIS deles, na sequencia das perguntas da disciplina.
ORDEM_FASE2 = [
    "linha_do_tempo",     # abertura do bloco: o que mudou
    "contrato_campos",    # dados novos
    "contrato_eda",       # EDA de contrato
    "contrato_h6b",       # EDA de contrato (H6b)
    "curadoria",          # curadoria de variaveis (02/09)
    "alvos",              # desenho da modelagem: Y1 e Y2
    "finais",             # item 12 - resultados finais
    "implicacoes",        # item 13
    "limitacoes_f2",      # item 14
    "futuros",            # item 15
    "conclusoes",         # item 16
]
mapa = dict(novos)
for chave in ORDEM_FASE2:
    el = id_do_slide(mapa[chave])
    sldIdLst.remove(el)
    sldIdLst.append(el)
print(f"bloco da Fase 2 anexado ao final: {len(ORDEM_FASE2)} slides")

def slide_dos_gates():
    """Localiza o slide de Gates pelo conteudo — a posicao muda quando o bloco da
    Fase 2 e anexado ao final."""
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.has_text_frame and "Entregas do Projeto" in sh.text_frame.text:
                return sl
    raise ValueError("slide de Gates nao encontrado")


# ============================================================
# 6) Gates: fechar Gate 4 e registrar o contrato
# ============================================================
gates_sl = slide_dos_gates()
for sh in gates_sl.shapes:
    if not sh.has_text_frame:
        continue
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            if "Resultados promissores" in r.text:
                r.text = (f"Y1 (custo anual): R² {br(_r2_dec)} pelo caminho decomposto e "
                          f"{br(best_expl['r2'])} explicativo. Y2 (nº de OS): R² {br(_y2['r2'])}. Teste 2025.")
            elif "Fatores: n" in r.text and "VMRS" in r.text:
                r.text = "Fatores: refrigeração, histórico defasado, exposição acumulada e subtipo."
            elif "Verificar estabilidade por ano de teste" in r.text:
                r.text = ("Contrato incorporado e testado: o TIPO importa (4ª variável mais "
                          f"importante); a duração, não ({float(cp['delta_r2_contrato']):+.4f} de R² no bloco).")


# ============================================================
# 7) Correcoes de revisao (2026-08-16)
# ============================================================
# Aplicadas sobre TODOS os slides, inclusive os acrescentados nesta fase.

CORRECOES_TEXTO = [
    # -- afirmacao que contradiz o resultado final: a segmentacao nao se sustentou --
    ("→ tratar a frota por grupos (refrigerada × seca), como no MERF, melhora a previsão.",
     "→ a refrigeração é a variável nº 1 do modelo; ela entra como atributo, não como "
     "critério de separação da frota."),
    # -- metodologia superada (slide oculto 'Desenho do estudo') --
    ("População: todo o custo interno; sem filtro MAINT, pois o dado de contrato não existe na fonte única.",
     "População: a frota completa (47.715 carreta-anos). O regime de manutenção entra como "
     "variável do modelo; o recorte por MAINT foi testado e mantido apenas como comparação."),
    ("Universo de variáveis: 25 candidatas da fonte — atributos do ativo, uso/quilometragem, geografia e histórico defasado.",
     "Universo de variáveis: fonte única com 29 colunas — atributos do ativo, uso/quilometragem, "
     "geografia, histórico defasado e contrato (H6 avaliada nesta fase)."),
    # -- moeda errada e mistura de limite do grafico com maximo dos dados --
    ("Distribuição fortemente assimétrica à direita — a maioria dos veículos custa pouco (mediana ~R$ 1.000) "
     "e uma cauda longa de outliers chega a R$ 12.000.",
     "Distribuição fortemente assimétrica à direita — a maioria das carretas custa pouco "
     "(mediana CAD 812/ano). O eixo do gráfico é cortado em ~CAD 12.000 (p99 = 11.804), "
     "mas a cauda real vai até CAD 62.231."),
    # -- KNN faltando na lista de tecnicas --
    ("Machine Learning: árvore de decisão, Random Forest, Gradient Boosting.",
     "Machine Learning: árvore de decisão, Random Forest, Gradient Boosting e KNN."),
    # -- afirmacao forte demais sobre zero-inflacao --
    ("O grão anual eliminou a zero-inflação que travava as formulações anteriores "
     f"({sv('share_y_zero_pct')}% de carreta-anos com custo zero).",
     f"O grão anual reduziu a zero-inflação a {sv('share_y_zero_pct')}%, deixando-a residual para a modelagem."),
]

# -- numeracao das secoes: sequencia unica e sem duplicidade --
RENUMERA = {
    # unica correcao no bloco de agosto: duplicidade de secao no referencial
    "12 · REFERENCIAL TEÓRICO": "11 · REFERENCIAL TEÓRICO",
}

def substituir_no_paragrafo(par, de, para):
    """Substitui texto que pode estar dividido em varios runs, preservando o formato
    do primeiro run do paragrafo."""
    inteiro = "".join(r.text for r in par.runs)
    if de not in inteiro:
        return False
    novo_txt = inteiro.replace(de, para)
    if not par.runs:
        return False
    par.runs[0].text = novo_txt
    for r in par.runs[1:]:
        r.text = ""
    return True


n_corr = 0
for sl in prs.slides:
    for sh in sl.shapes:
        if not sh.has_text_frame:
            continue
        bruto = sh.text_frame.text.strip()
        if bruto in RENUMERA:
            for p_ in sh.text_frame.paragraphs:
                if substituir_no_paragrafo(p_, bruto, RENUMERA[bruto]):
                    n_corr += 1
            continue
        for p_ in sh.text_frame.paragraphs:
            for a, b in CORRECOES_TEXTO:
                if substituir_no_paragrafo(p_, a, b):
                    n_corr += 1
print(f"correcoes de revisao aplicadas: {n_corr}")

# -- Gates: status coerente (Gate 4 e Gate 5 concluidos; so o Gate 6 fica pendente) --
gates_sl = slide_dos_gates()
STATUS = {(3.88, 11.36): "Concluída",   # Gate 4 - Modelo anual de custo
          (5.29, 5.18): "Concluída"}    # Gate 5 - Validacao e robustez
for sh in gates_sl.shapes:
    if not sh.has_text_frame:
        continue
    if sh.text_frame.text.strip() not in ("Próximos passos", "Concluída"):
        continue
    top_in = round(sh.top / 914400, 2)
    left_in = round(sh.left / 914400, 2)
    for (t_, l_), novo_status in STATUS.items():
        if abs(top_in - t_) < 0.06 and abs(left_in - l_) < 0.06:
            for p_ in sh.text_frame.paragraphs:
                for r_ in p_.runs:
                    r_.text = novo_status
            print(f"  gate em (top={t_}, left={l_}) -> {novo_status}")

# ordem de leitura: no XML o bloco do Gate 6 vem antes do Gate 5. Reposiciona os
# elementos para que a ordem do arquivo acompanhe a ordem visual (esquerda -> direita).
spTree = gates_sl.shapes._spTree
blocos = [sh for sh in gates_sl.shapes if sh.has_text_frame and sh.top is not None
          and abs(sh.top / 914400 - 5.28) < 0.45]
if blocos:
    for sh in sorted(blocos, key=lambda s: (s.left, s.top)):
        el = sh._element
        spTree.remove(el)
        spTree.append(el)
    print("  ordem de leitura dos Gates 5/6 normalizada")


# ============================================================
# 8) Correcoes de PRECISAO (revisao 2026-08-16)
# ============================================================
PRECISAO = [
    ("Critérios: ranking, VIF, redundância, vazamento temporal e coerência de domínio.",
     "Seleção conforme a curadoria de variáveis do Grupo (02/set), confrontada com a evidência: "
     "ranking de associação, VIF, redundância, vazamento temporal e coerência de domínio."),

    # -- 1. definicao da populacao: escopo do custo != populacao de modelagem --
    ("Escopo: todo o custo interno absorvido pela empresa (preventiva + corretiva).",
     "Escopo do custo: todo o custo interno absorvido pela empresa (preventiva + corretiva). "
     "População de modelagem: a frota completa — 47.715 carreta-anos. O tipo de contrato "
     "entra como variável do modelo, e não como recorte da amostra."),

    # -- 2. janela temporal: 6 anos de dados, 5 modelaveis --
    ("Fonte única   — CSV consolidado fato_wo_ml (217.217 OS · 9.585 carretas únicas).",
     "Fonte única — CSV consolidado fato_wo_ml (217.217 OS · 9.585 carretas · 2020–2025)."),
    ("Base carreta × ano   —  define o alvo Y = custo anual de manutenção por carreta (47.715 linhas).",
     "Base carreta × ano — alvos Y1 = custo anual por carreta e Y2 = nº de OS no ano (47.715 linhas). "
     "São 6 anos de dados, mas só 5 modeláveis: 2020 não tem ano anterior para as variáveis defasadas."),
    ("Modelagem — split temporal   —  treino 2020–2024 / teste 2025; lineares, árvores e ensembles.",
     "Modelagem — split temporal — treino 2020–2024 / teste 2025; validação adicional com 2023 e 2024 "
     "como anos de teste; lineares, árvores e ensembles."),

    # -- 3. importancia por permutacao: embaralhar != remover --
    ("Refrigeração é o que mais pesa: carretas refrigeradas custam mais e são a variável nº 1 do modelo "
     "(removê-la derruba o R² em 0,17).",
     f"Refrigeração é a variável nº 1 do modelo: embaralhar seus valores reduz o R² em {br(_imp_ref, 2)} "
     "(importância por permutação) — mais que todas as outras somadas. Isso mede a dependência do "
     "modelo ajustado, não equivale a retreinar sem a variável."),

    # -- 8. inferencias atribuidas a literatura --
    ("→ custo × operação é não linear, como no Super Learner de Katreddi (R²≈97%).",
     "→ a relação custo × operação é não linear, como reporta Katreddi (2023). O R² de 97% daquele "
     "estudo não é comparável ao nosso: outra frota, outro alvo e outra unidade."),
    ("→ confirma Sun (histórico prevê custo); a idade sozinha não manda.",
     "→ resultado consistente com Sun (2024), que usa o histórico para prever custo; aqui a idade, "
     "isolada, tem associação fraca."),
    ("→ número honesto; avaliação temporal como defende a literatura.",
     "→ a queda mostra quanto do R² explicativo vinha de variáveis do próprio ano; a avaliação "
     "temporal segue a prática recomendada na literatura."),

    # -- 5. conclusao sobre idade: mediacao nao foi testada --
    ("A idade isolada NÃO explica o custo: atua por meio do histórico e do uso.",
     "A idade, isoladamente, tem associação fraca com o custo (ρ = 0,03) e baixa importância no modelo. "
     "Não testamos se o efeito da idade é mediado por histórico e uso — é hipótese, não resultado."),

    # -- 4 e 6. MAE e alegacao de aplicacao operacional --
    ("O modelo é utilizável para orçamento e priorização, desde que a margem de erro seja declarada.",
     "O erro médio (MAE) é de CAD 1.059/ano, cerca de 52% do custo médio do ano de teste: o modelo "
     "serve para PRIORIZAR carretas e para provisionar no agregado da frota, não para estimar o "
     "valor de um ativo isolado."),
    (f"Erro médio do modelo recomendado: CAD $ {best_pred['mae']}/ano por carreta — margem a declarar ao usar a "
     "estimativa em orçamento.",
     f"MAE = CAD {float(_dec['mae']):.0f}/ano por carreta — é o erro MÉDIO, não um limite: metade dos casos erra "
     "menos e a cauda erra muito mais. Equivale a ~52% do custo médio do ano de teste (CAD 2.034)."),
    ("Aplicação ao planejamento   —  apoia o controle dos custos internos e a elaboração de orçamentos "
     "mais aderentes à idade, às especificações e à aplicação de cada equipamento da frota.",
     "Aplicação ao planejamento — apoia a priorização de carretas e o provisionamento agregado por "
     "perfil de frota; o erro por ativo isolado é alto demais para orçamento individual."),
    ("Deployment e atualização contínua   —  disponibiliza os scripts de ML para integração com APIs "
     "que alimentem o pipeline trimestralmente, permitindo revalidar e atualizar os resultados do modelo.",
     "Deployment e atualização — o pipeline é reprodutível ponta a ponta pelos notebooks; a integração "
     "com APIs e a atualização periódica são etapa futura, ainda não implementada."),

    # -- 7. Gates: descrever o que foi de fato entregue --
    ("Resultados promissores, superiores aos modelos lineares no teste temporal.",
     "Y1: R² 0,471 pelo caminho decomposto (0,442 direto) e 0,578 explicativo no teste 2025. Y2: R² 0,608."),
    ("Fatores: nº de OS, VMRS, custo acumulado, subtipo, km e região.",
     "Fatores dominantes: refrigeração, histórico defasado, exposição acumulada e subtipo."),
    ("Verificar estabilidade por ano de teste, região, tipo de carreta e perfil de manutenção.",
     "FEITO: estabilidade verificada com 3 anos de teste (2023–2025); contrato testado — o TIPO importa, a duração não."),
    ("Testar sensibilidade a outliers, seleção de variáveis e hiperparâmetros.",
     "PENDENTE: sensibilidade a outliers, seleção de variáveis e hiperparâmetros."),
    ("Avaliar extensões: MERF, modelos hierárquicos e abordagens zero-infladas.",
     "PENDENTE: MERF, modelos hierárquicos e a previsão de 2026."),
]

n_prec = 0
for sl in prs.slides:
    for sh in sl.shapes:
        if not sh.has_text_frame:
            continue
        for p_ in sh.text_frame.paragraphs:
            for a, b in PRECISAO:
                if substituir_no_paragrafo(p_, a, b):
                    n_prec += 1
print(f"correcoes de precisao aplicadas: {n_prec}")

# Gate 5 nao esta concluido: parte do escopo segue pendente
for sh in slide_dos_gates().shapes:
    if not sh.has_text_frame:
        continue
    if sh.text_frame.text.strip() != "Concluída":
        continue
    if abs(sh.top / 914400 - 5.29) < 0.06 and abs(sh.left / 914400 - 5.18) < 0.06:
        for p_ in sh.text_frame.paragraphs:
            substituir_no_paragrafo(p_, "Concluída", "Parcial")
        print("  Gate 5 -> Parcial (validação feita; robustez pendente)")

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"OK: {OUT.name} - {len(prs.slides._sldIdLst)} slides")
